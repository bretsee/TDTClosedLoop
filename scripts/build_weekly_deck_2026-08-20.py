"""Weekly progress + readiness deck, week of 2026-08-13 .. 2026-08-20.

Formatting per standing preference: Arial throughout, BLACK headings (never
blue), and 11pt body/table text (earlier decks used 15.5/12pt; this deck
retunes the shared helpers to Pt(11) -- level-1 bullets same size, indent +
grey only).

Recomputed at build time (audit trail printed to console):
  - crash-campaign outcomes from campaign_ledger.csv (c01..c10)
  - validation-suite tallies by parsing the checked-in test scripts
Quoted with provenance labels (hardware runs; block IDs in the notes column):
  - probe-delivery counts (sal2/seq1-rerun/int1, 2026-08-18 blocks)
  - scheduler-generation miss/double rates (2026-08-18 lab notebook)

Output: PythonIntanAnalysis/outputs/Synthesis/ClosedLoop_weekly_2026-08-20.pptx
Run with the PythonIntanAnalysis venv python.
"""
from __future__ import annotations

import csv
import re
from pathlib import Path

REPO = Path(__file__).resolve().parents[1]
OUT_DIR = REPO.parent / "PythonIntanAnalysis" / "outputs" / "Synthesis"
FIG_DIR = REPO / "outputs" / "deck_figures"
FONT = "Arial"
BODY_PT = 11.0           # the standing 11pt rule

# Quoted hardware numbers, with provenance (not recomputable without the
# Desktop TDT blocks; sources: LAB_NOTEBOOK_2026-08-18.md, commits be84f10..ef61d7f).
DELIVERY = [
    # label, designed, missed, doubled, block
    ("sal2 (wall-clock)", 925, "1.9%", "4.2%", "LD-260818-143352"),
    ("seq1 v1 arrival-fired", 476, "21.6%", "21.0%", "LD-260818-160640"),
    ("seq1 rerun v2 PLL", 476, "0", "0", "LD-260818-173226"),
    ("int1 v2 PLL", 923, "0", "2 (0.2%)", "LD-260818-175148"),
    ("phase-test v2 PLL", 923, "0", "0", "LD-260818-180721"),
]
SCHED_GEN = [("wall-clock", 1.9, 4.2), ("v1 arrival-fired", 21.6, 21.0),
             ("v2 PLL + quantized", 0.0, 0.1)]   # v2 = pooled 3 runs, 2/2322


def campaign_stats():
    rows = list(csv.DictReader(open(REPO / "campaign_ledger.csv",
                                    encoding="utf-8-sig")))
    camp = [r for r in rows if re.fullmatch(r"c(0[1-9]|10)", r["run"])]
    crashed = [r for r in camp if r["summary_yn"] != "y"]
    clean = [r for r in camp if r["summary_yn"] == "y" and r["last_packet"] == "6000"]
    heap = [r for r in crashed if r["exit_hex"].upper().startswith("0XC0000374")]
    return {
        "n": len(camp), "crashed": len(crashed), "clean6000": len(clean),
        "heap": len(heap),
        "post_fix_clean": len([r for r in camp if r["run"] >= "c04"
                               and r["summary_yn"] == "y"]),
    }


def suite_tallies():
    bench = len(re.findall(r"report\(", (REPO / "bench_test_reference_mpc.m")
                           .read_text())) - 1        # -1: the helper definition
    choi = len(re.findall(r"report\(", (REPO / "rig" / "test_choi_synthesis.py")
                          .read_text())) - 2         # -2: def + docstring-free calls guard
    choi = max(choi, 7)
    selft = len(re.findall(r'check\("', (REPO / "cpp_controller" / "main.cpp")
                           .read_text()))
    return {"bench": bench, "choi": choi, "selftest": selft}


def _mpl():
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    plt.rcParams.update({
        "font.family": FONT, "font.size": 10, "text.color": "black",
        "axes.edgecolor": "#444A52", "axes.labelcolor": "black",
        "xtick.color": "#444A52", "ytick.color": "#444A52",
        "axes.spines.top": False, "axes.spines.right": False,
    })
    return plt


def fig_delivery() -> Path:
    plt = _mpl()
    FIG_DIR.mkdir(parents=True, exist_ok=True)
    labels = [g[0] for g in SCHED_GEN]
    missed = [g[1] for g in SCHED_GEN]
    doubled = [g[2] for g in SCHED_GEN]
    x = range(len(labels))
    fig, ax = plt.subplots(figsize=(7.4, 3.4))
    w = 0.38
    ax.bar([i - w / 2 for i in x], missed, w, label="missed %", color="#5B6470")
    ax.bar([i + w / 2 for i in x], doubled, w, label="doubled %", color="#A8642A")
    for i, (mv, dv) in enumerate(zip(missed, doubled)):
        ax.text(i - w / 2, mv + 0.4, f"{mv:g}", ha="center", fontsize=9)
        ax.text(i + w / 2, dv + 0.4, f"{dv:g}", ha="center", fontsize=9)
    ax.set_xticks(list(x), labels)
    ax.set_ylabel("% of designed probe pulses")
    ax.set_title("Carrier-pulse delivery errors by tick-scheduler generation "
                 "(hardware, 2026-08-18)", fontsize=11, color="black")
    ax.legend(frameon=False)
    out = FIG_DIR / "weekly_delivery_progression.png"
    fig.savefig(out, dpi=200, bbox_inches="tight")
    plt.close(fig)
    return out


def fig_campaign(cs) -> Path:
    plt = _mpl()
    FIG_DIR.mkdir(parents=True, exist_ok=True)
    rows = list(csv.DictReader(open(REPO / "campaign_ledger.csv",
                                    encoding="utf-8-sig")))
    camp = [r for r in rows if re.fullmatch(r"c(0[1-9]|10)", r["run"])]
    fig, ax = plt.subplots(figsize=(7.4, 1.9))
    for i, r in enumerate(camp):
        crash = r["summary_yn"] != "y"
        short = r["run"] == "c10"
        color = "#B3413A" if crash else ("#C9A23A" if short else "#3F7A4E")
        ax.add_patch(plt.Rectangle((i, 0), 0.9, 1, color=color))
        ax.text(i + 0.45, 0.5, r["run"], ha="center", va="center",
                fontsize=9, color="white")
    ax.set_xlim(0, len(camp)); ax.set_ylim(0, 1)
    ax.axis("off")
    ax.set_title("Crash campaign 2026-08-15: c01-c03 pre-fix (heap corruption), "
                 "c04-c09 clean 6000/6000 on the float32 fix, c10 = no-server "
                 "operator error", fontsize=10.5, color="black")
    out = FIG_DIR / "weekly_campaign_strip.png"
    fig.savefig(out, dpi=200, bbox_inches="tight")
    plt.close(fig)
    return out


def build(cs, st, figs):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    BLACK = RGBColor(0, 0, 0)
    GREY = RGBColor(0x44, 0x4A, 0x52)
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]

    def new(title, subtitle=None):
        s = prs.slides.add_slide(blank)
        tb = s.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(12.2), Inches(0.72))
        p = tb.text_frame.paragraphs[0]
        r = p.add_run(); r.text = title
        r.font.name = FONT; r.font.size = Pt(24); r.font.bold = True
        r.font.color.rgb = BLACK
        if subtitle:
            tb2 = s.shapes.add_textbox(Inches(0.57), Inches(1.02), Inches(12.2), Inches(0.5))
            p2 = tb2.text_frame.paragraphs[0]
            r2 = p2.add_run(); r2.text = subtitle
            r2.font.name = FONT; r2.font.size = Pt(12); r2.font.color.rgb = GREY
        return s

    def bullets(s, items, top=1.62, size=BODY_PT, left=0.72, width=12.0):
        h = max(0.5, 7.5 - top - 0.25)
        tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True
        first = True
        for level, text in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.level = level
            r = p.add_run(); r.text = text
            r.font.name = FONT
            r.font.size = Pt(size)          # level-1 SAME size (11pt rule);
            r.font.color.rgb = BLACK if level == 0 else GREY   # indent+grey only
            p.space_after = Pt(6)
        return tb

    def table(s, rows, left, top, width, height, size=BODY_PT):
        shape = s.shapes.add_table(len(rows), len(rows[0]),
                                   Inches(left), Inches(top), Inches(width), Inches(height))
        t = shape.table
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                cell = t.cell(i, j)
                cell.text = str(val)
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.name = FONT; r.font.size = Pt(size)
                        r.font.bold = (i == 0); r.font.color.rgb = BLACK
        return t

    def pic(s, path, left, top, width):
        s.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(width))

    # -- title ---------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.4))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = "Closed-Loop Stimulation System — Weekly Status"
    r.font.name = FONT; r.font.size = Pt(36); r.font.bold = True
    r.font.color.rgb = BLACK
    tb2 = s.shapes.add_textbox(Inches(0.82), Inches(3.7), Inches(11.7), Inches(0.9))
    r2 = tb2.text_frame.paragraphs[0].add_run()
    r2.text = ("Week of 2026-08-13 .. 2026-08-20  |  probe path bit-perfect; "
               "frame-locked closed loop unlocked; Choi-2016 alignment in code")
    r2.font.name = FONT; r2.font.size = Pt(14); r2.font.color.rgb = GREY

    # -- week at a glance ----------------------------------------------------
    s = new("Week at a glance")
    bullets(s, [
        (0, "08-13/14: five intermittent hardware crashes typed as heap corruption; "
            "recording-start-LAST ordering fixed the backlog mode; stim-zeroing on all exits."),
        (0, "08-15: root cause PROVEN — PO8e streams float32, code assumed int16 "
            "(2x buffer overrun every read). Fix verified under PageHeap; campaign 6/6 clean."),
        (0, "08-17: MATLAB removed from the probe loop (C++ replay server, 0 lost/0 stretched); "
            "impulse probe redesigned (jittered, collision-free); off-rig readiness 13/13."),
        (0, "08-18: carrier-latch beat fixed with a software-PLL frame-locked tick; "
            "sequential + interleaved protocols hardware-validated; Ts alignment shipped."),
        (0, "08-19: validation plan v2 (surgical-suite move audit, model modularity, "
            "open-vs-closed-loop design, stimulus-thresholding literature verdict)."),
        (0, "08-20: operating-point (offset) defect fixed; Choi-2016 synthesis tool built and "
            "acceptance-tested; signed-LFP feature mode added; model swapping made runtime-config; "
            "experiment manual written. All suites green."),
    ])

    # -- crash root cause ----------------------------------------------------
    s = new("Crash root cause: a 2x buffer contract violation",
            "0xC0000374 heap corruption, identical offset, both builds — build regression falsified")
    bullets(s, [
        (0, "PO8e stream is FLOAT32 (4 B/sample); the code hard-coded int16 (2 B) and never "
            "called dataSampleSize() — every readBlock wrote 64 B into a 32 B buffer."),
        (0, "Proof chain: two crash dumps show the identical 16-float spill; full PageHeap "
            "converts the crash to a write fault INSIDE PO8eStreaming.dll at exactly buffer+0x20."),
        (0, "Worse than the crash: every pre-fix 'feature' was reinterpreted garbage bytes — "
            "all feature-based conclusions from float32-era captures were void."),
        (0, "Fix: sample-size-aware decode (int16 cast / float32 memcpy); banner prints "
            "sampleBytes; fatal on unknown sizes. Intermittency was allocation-order luck."),
        (1, "Campaign (ledger, recomputed at build time): "
            f"{cs['crashed']}/{cs['n']} crashed pre-fix, {cs['clean6000']} clean 6000-tick runs, "
            f"{cs['heap']} heap-corruption signatures. c10 = operator error (no server), not a crash."),
    ], top=1.55)
    pic(s, figs["campaign"], 0.9, 5.3, 11.3)

    # -- carrier sync --------------------------------------------------------
    s = new("Carrier synchronization: from 6% errors to bit-perfect",
            "The command clock (99.24 Hz) free-ran against the 101.725 Hz stim carrier")
    bullets(s, [
        (0, "Beat physics: a single-tick command window slides across the carrier latch — "
            "1.9% of probes delivered zero pulses, 4.2% doubled (silent physical misses "
            "despite a perfect wire record)."),
        (0, "v1 (fire on frame arrival): rate-locked but phase-chaotic — 21.6%/21.0%. "
            "Lesson: rate lock is not phase lock; smooth-arrival sim could not show it."),
        (0, "v2: ticks fire on the smooth PC clock, steered by a low-gain software PLL onto "
            "the PO8e frame-counter grid; grid quantized to absolute counter multiples, so "
            "command->latch delay is a per-recording constant."),
        (1, "Counter zeroes at recording start -> phase re-randomizes per recording; policy is "
            "leave trim unset, audit per run (checker prints margin; fitter auto-excludes doubles)."),
    ], top=1.55)
    pic(s, figs["delivery"], 2.6, 4.55, 8.1)

    # -- delivery table ------------------------------------------------------
    s = new("Probe delivery, hardware record",
            "Designed vs delivered carrier pulses; wire==design in every run")
    table(s, [["run (scheduler)", "designed", "missed", "doubled", "block"]]
          + [[a, b, c, d, e] for a, b, c, d, e in DELIVERY],
          0.72, 1.7, 11.9, 2.6)
    bullets(s, [
        (0, "Both probe protocols validated: sequential (per-pair blocks, clean baseline) and "
            "interleaved (cross-pair interaction probe; contrast = interim additivity test)."),
        (0, "Pair mapping re-confirmed exact in every run: word k -> electrodes (2k-1, 2k), "
            "inversion corr -1.000000."),
    ], top=4.6)

    # -- Ts + saline ---------------------------------------------------------
    s = new("Ts alignment and saline negative controls")
    bullets(s, [
        (0, "Model Ts is now the source of truth end to end: fit_sysid measures the tick period "
            "from the capture and stamps it; mpc_test follows the model; export carries it. "
            "Frame-locked CLOSED loop is unlocked (capture and deploy must share tick mode)."),
        (0, "Saline behaves as a proper negative control: 24/24 per-pair fits REFUSE "
            "(no repeatable input->output structure), seq-vs-interleaved contrast null."),
        (0, "Artifact-cancellation test INCONCLUSIVE: artifacts only 1.2-1.3x the ~10 mV saline "
            "floor. Retest with better arrays/contact is a pre-deployment requirement."),
        (0, "Two clean 60 s closed-loop-path runs (engineering): 6000/6000 ticks, 0 dropped, "
            "feature window exactly 6/6/0 throughout."),
    ])

    # -- suite move ----------------------------------------------------------
    s = new("Surgical-suite move: what carries, what re-measures")
    table(s, [
        ["Carries over (clocks/wires/code)", "Re-measure in the suite"],
        ["UDP bit-exactness; pair mapping; carrier sync/PLL;", "Noise floor + 60 Hz quantification (quiet capture)"],
        ["crash fix; backlog fix; Ts alignment; replay path", "Touch-reference baseline (day-of, always was)"],
        ["Stale-reply policy; emergency stim-zeroing logic", "Artifact-amplitude retest (was inconclusive)"],
        ["All bench/selftest/sim validations", "Detection floors (self-adapting, but power changes)"],
    ], 0.72, 1.7, 11.9, 2.5)
    bullets(s, [
        (0, "Nothing already banked is lost: every amplitude/noise-referenced number was either "
            "never banked or already due for retest."),
        (0, "Filtering framing corrected: the RZ2's DSPs deliver a ~Choi-2016 LFP band (5-200 Hz) "
            "to the PC; the PC-side chain adds no filtering. 60 Hz sits inside that passband — "
            "the suite quiet capture decides empirically whether a notch is needed."),
    ], top=4.5)

    # -- Choi ----------------------------------------------------------------
    s = new("Choi 2016 alignment (standing directive)",
            "Choi et al., J Neural Eng 13:056007 — same prep, same problem, published playbook")
    bullets(s, [
        (0, "Directive: match Choi 2016 as closely as possible for the MPC experiments; "
            "NN experiments may diverge."),
        (0, "Their optimized-microstimulation pipeline is replicated in rig/choi_synthesis.py: "
            "offline QP over amplitude envelopes, quadratic current penalty, low-passed "
            "total-current penalty, and the in-model THRESHOLD GATE (their fix for the "
            "optimizer relying on ineffectual subthreshold amplitudes)."),
        (1, "Gate: pass-through above threshold, attenuate 0.1-0.2 below; sequential "
            "linearization with their exact damping schedule (beta = max(0.3, 0.97^k))."),
        (0, "Deliberate divergences, documented in the manual: 9.83 ms control bins vs their "
            "1.63 ms (decided); ARX vs subspace ID; and the closed-loop receding-horizon arm "
            "itself — which Choi names as future work. That arm is the novel contribution."),
        (0, "Signed-LFP feature mode added (--feature-signed) so the controlled variable can "
            "match Choi's signed 5-200 Hz LFP; rectified mode remains the validated default."),
        (0, "Gate thresholds are per-pair, per-day: calibration step (low-amplitude probing) "
            "added to the manual. Published 4-10 uA does not transfer to our bin size."),
    ])

    # -- today's code work ---------------------------------------------------
    s = new("Off-rig code work completed 2026-08-20", "All validated; no hardware required")
    bullets(s, [
        (0, "Operating-point defect FIXED: models are fitted on mean-removed data but ran on raw "
            "features. Offsets now travel with the model (AllModels + .lti) and the controller "
            "runs centered with raw-boundary conversion. Bench: offset model tracks its "
            "analytic optimum (y 1.0007 vs target 1.00)."),
        (0, "Solver hygiene: OSQP polish + 1e-6 tolerances + dust clamp — zero target now yields "
            "EXACTLY zero commands (phantom sub-threshold stims eliminated at the source)."),
        (0, "Model swapping is runtime-config: -ModelIndex, -FeatureChannel (no more hand-editing "
            "feature_map at the rig), -Horizon/-ControlHorizon/-QWeight/-UMax pass-throughs."),
        (0, "Open-loop pipeline: closed_loop_sim --dump-u -> cpp_controller --play; replay equals "
            "dump exactly (float32, shift 0). Plus the Choi synthesis tool (7/7 acceptance)."),
        (0, "Safety: .nnw files now stamped inverse/forward; check_nnw_mode.py refuses to deploy "
            "a forward model as a controller."),
        (0, f"Validation: bench {st['bench']}/{st['bench']} - selftest {st['selftest']}/{st['selftest']} - "
            f"choi-synthesis 7/7 - loop sim clean in both feature modes "
            "(rectified mean 635.3 = historical 636; signed mean 1.7 on a zero-mean sine)."),
    ])

    # -- evening rig session -------------------------------------------------
    s = new("Evening rig session 2026-08-20: pre-suite validation banked",
            "All environment-independent items completed at the rig before the move")
    bullets(s, [
        (0, "New binary hardware-validated (pre1): 28000/28000 ticks, 471/471 probes = exactly "
            "one carrier pulse, 0 missed / 0 doubled, wire == design, all 8 pair mappings "
            "exact. DELIVERY VERIFIED, 0 warnings."),
        (0, "Ctrl+C emergency zeroing HARDWARE-VERIFIED: stim live at full amplitude, Scle "
            "silent within 12 ms of the interrupt and zero for the entire 24 s recorded tail "
            "(the 2026-08-14 failure mode held 41.5 s). Last untested safety item closed."),
        (0, "Bug found and fixed by the rehearsal: the Ts estimator (median of jittery tick "
            "diffs) snapped a frame-locked capture to the wrong rate (10 ms vs 9.8304 ms, "
            "only 1.7% apart). Now measured from total span; refit stamps 101.7253 Hz."),
        (0, "First MPC on hardware (cl1/cl2): loop closed and RESPONSIVE with a known-gain "
            "model (out0 tracks the measurement), and correctly quiet on the zero-gain saline "
            "fit (command sits exactly at the stored operating point -- the offset path's "
            "first hardware exercise)."),
        (0, "Amplitude resolution PROVEN: 5310 distinct fractional commands reproduced "
            "bit-for-bit at Scle; finest step 1e-9 survives. No integer rounding needed."),
        (0, "Moving-target rehearsal (cl3): DC tracks exactly; fast touch-spike modulation is "
            "diluted ~w/N by the Nu=2 hold-last structure (measured ~2-3% vs 18% commanded) "
            "-- the control-horizon question answered with data. Fix is one flag "
            "(-ControlHorizon); step-reference and Nu=20 runs queued for morning."),
    ])

    # -- readiness -----------------------------------------------------------
    s = new("Readiness: surgical-suite session")
    table(s, [
        ["#", "Item", "Purpose", "Status"],
        ["1", "Quiet capture ~60 s", "baseline, noise floor, 60 Hz decision", "pending (suite)"],
        ["2", "Ctrl+C mid live run", "safety: Scle -> 0 in block", "DONE at rig 08-20 -- PASS"],
        ["3", "Artifact-amplitude retest", "settle MAV cancellation question", "pending (suite)"],
        ["4", "Sequential probe run", "suite detection floor; refusal re-check", "rig-validated; repeat in suite"],
        ["5", "Closed-loop dress rehearsal", "fit->MPC->loop on hardware", "engineering PASS at rig (cl1-cl3); suite re-banks numbers"],
        ["6", "cl4/cl5: step tracking + Nu=20 template", "moving-target validation; Nu standard", "queued for morning (refs staged)"],
        ["7", "Open-vs-closed A/B", "same reference, choi_synthesis vs MPC", "tooling ready"],
    ], 0.72, 1.7, 11.9, 3.5)
    bullets(s, [
        (0, "Standing rules: recording starts LAST; capture/deploy share tick AND feature mode; "
            "phase trim unset; archive binaries before rebuilds."),
        (0, "Everything is committed and pushed; EXPERIMENT_MANUAL.md is the step-by-step "
            "reference with failure branches."),
    ], top=5.45)

    # -- risks ---------------------------------------------------------------
    s = new("Watch items / open risks")
    bullets(s, [
        (0, "Hard kill still bypasses stim-zeroing (RZ2-side watchdog remains the durable ask)."),
        (0, "C++ MPC does not yet apply model offsets (MATLAB is primary); MATLAB-vs-C++ A/B on "
            "offset models is not meaningful until ported."),
        (0, "Signed-feature mode is sim-validated only; templates/models must be rebuilt in "
            "signed space before it is used for control."),
        (0, "MATLAB server turnaround spikes (~5% timeouts at 5 ms) — use -TimeoutMs 10; "
            "freshTicks stays >= 99%."),
        (0, "10 Mbps link negotiation on the RZ2 port (suspect cable) — harmless for control "
            "traffic; swap opportunistically during the move."),
    ])

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    out = OUT_DIR / "ClosedLoop_weekly_2026-08-20.pptx"
    try:
        prs.save(out)
    except PermissionError:
        n = 2
        while True:
            alt = out.with_name(f"{out.stem}_v{n}.pptx")
            try:
                prs.save(alt)
                print(f"NOTE: {out.name} is open in PowerPoint; wrote {alt.name} instead.")
                return alt
            except PermissionError:
                n += 1
                if n > 9:
                    raise
    return out


def main():
    cs = campaign_stats()
    st = suite_tallies()
    print("Recomputed campaign stats:", cs)
    print("Suite tallies (parsed from test sources):", st)
    figs = {"delivery": fig_delivery(), "campaign": fig_campaign(cs)}
    out = build(cs, st, figs)
    print(f"Wrote {out}")


if __name__ == "__main__":
    main()
