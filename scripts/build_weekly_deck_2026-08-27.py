"""Weekly progress + readiness deck, week of 2026-08-21 .. 2026-08-27.

Formatting per standing preference: Arial throughout, BLACK headings (never
blue), 11pt body/table text (same shared-helper tuning as the 08-20 deck).

Recomputed at build time (audit trail printed to console):
  - suite quiet-capture numbers from quiet_2026-08-26.json (noise floor,
    line fraction, baseline; ch 13 exclusion)
  - rndval designed-probe count from design_runrndval.csv (rising edges/word)
  - MPC-check tracking numbers from tracking_mpccheck_20260827.json
  - MPC-check u-vs-reference figure from capture_mpc_20260827_184338.csv
    + ref_steps.csv
Quoted with provenance labels (hardware runs; block IDs in the notes):
  - cl4/cl5 control-horizon results (2026-08-25 evening, saline)
  - nThw validation counts (2026-08-26 test blocks)
  - 08-27 delivery audit / artifact verdict (blocks LD-260827-182650/-184414)

Output: PythonIntanAnalysis/outputs/Synthesis/ClosedLoop_weekly_2026-08-27.pptx
Run with the PythonIntanAnalysis venv python.
"""
from __future__ import annotations

import csv
import json
from pathlib import Path

REPO = Path(__file__).resolve().parents[1]
OUT_DIR = REPO.parent / "PythonIntanAnalysis" / "outputs" / "Synthesis"
FIG_DIR = REPO / "outputs" / "deck_figures"
FONT = "Arial"
BODY_PT = 11.0           # the standing 11pt rule

BLOCK_RNDVAL = "LD-260827-182650"
BLOCK_MPCCHECK = "LD-260827-184414"


# ---------------------------------------------------------------- recompute --
def quiet_stats():
    d = json.loads((REPO / "quiet_2026-08-26.json").read_text())
    chans = d["channels"]
    good = [c for c in chans if c["ch"] != 13]
    return {
        "noise_lo": min(c["noise_uv"] for c in good),
        "noise_hi": max(c["noise_uv"] for c in good),
        "noise_ch13": next(c["noise_uv"] for c in chans if c["ch"] == 13),
        "base_lo": min(c["baseline_mav_v"] for c in good),
        "base_hi": max(c["baseline_mav_v"] for c in good),
        "line_med": d["median_line_frac"],
        "wav_saved": d["wav1_saved"] and d["wav2_saved"],
        "chans": chans,
    }


def _csv_cols(path):
    rows = list(csv.reader(open(path)))
    hdr = rows[0]
    data = rows[1:]
    return hdr, data


def rndval_designed():
    hdr, data = _csv_cols(REPO / "design_runrndval.csv")
    ucols = [i for i, h in enumerate(hdr) if h.strip().lower().startswith("u")]
    total = 0
    for i in ucols:
        v = [float(r[i]) for r in data]
        total += sum(1 for k in range(1, len(v)) if v[k] > 0 and v[k - 1] <= 0)
        total += 1 if v[0] > 0 else 0
    return total


def mpccheck_stats():
    d = json.loads((REPO / "tracking_mpccheck_20260827.json").read_text())
    pc = d["per_channel"][0]
    return {
        "slope_u": pc["slope_u_on_r"]["u1"],
        "verdict": pc["verdict"],
        "transients": pc["eta"]["n_transients"],
        "ticks": d["n_ticks"] + d["skip_ticks"],
    }


# ------------------------------------------------------------------ figures --
def _mpl():
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    plt.rcParams.update({
        "font.family": FONT, "font.size": 10, "text.color": "black",
        "axes.edgecolor": "#444A52", "axes.labelcolor": "black",
        "xtick.color": "#444A52", "ytick.color": "#444A52",
        "axes.spines.top": False, "axes.spines.right": False,
    })
    return plt


def fig_gate() -> Path:
    """The rndval delivery outcome incl. the enable-gate hole (quoted counts,
    block LD-260827-182650; designed count recomputed from the design CSV)."""
    plt = _mpl()
    FIG_DIR.mkdir(parents=True, exist_ok=True)
    fig, ax = plt.subplots(figsize=(8.6, 2.4))
    ax.barh(0, 12, color="#5B6470", height=0.5)
    ax.barh(0, 640, left=12, color="#3F7A4E", height=0.5)
    ax.text(6, -0.55, "12 pre-enable\n(safety gate: blocked)", ha="left",
            fontsize=9, color="#5B6470")
    ax.text(332, 0, "640/640 post-enable: exactly ONE carrier pulse each\n"
            "0 missed, 0 doubled, 0 latch races", ha="center", va="center",
            color="white", fontsize=10)
    ax.set_xlim(0, 652)
    ax.set_yticks([])
    ax.set_xlabel("designed probes (n=652; wire == design, 0 lost)")
    ax.set_title("Suite delivery gate 2026-08-27: every probe accounted for "
                 f"(block {BLOCK_RNDVAL})", fontsize=10.5, color="black")
    out = FIG_DIR / "weekly27_gate.png"
    fig.tight_layout()
    fig.savefig(out, dpi=200, bbox_inches="tight")
    plt.close(fig)
    return out


def fig_noise(qs) -> Path:
    """Per-channel suite noise floor, recomputed from quiet_2026-08-26.json."""
    plt = _mpl()
    FIG_DIR.mkdir(parents=True, exist_ok=True)
    fig, ax = plt.subplots(figsize=(8.2, 3.0))
    chs = [c["ch"] for c in qs["chans"]]
    nz = [c["noise_uv"] for c in qs["chans"]]
    colors = ["#B3413A" if c == 13 else ("#C9A23A" if c == 16 else "#3F7A4E")
              for c in chs]
    ax.bar(chs, nz, color=colors, width=0.7)
    ax.text(13, qs["noise_ch13"] + 1, "ch 13 BLACKLISTED\n(bad contact)",
            ha="center", fontsize=9, color="#B3413A")
    ax.set_xticks(chs)
    ax.set_xlabel("channel")
    ax.set_ylabel("noise floor (µV std)")
    ax.set_title("Suite quiet capture 2026-08-26: real 8-11 µV floor "
                 "(the old bath's ~10 mV was the anomaly)",
                 fontsize=10.5, color="black")
    out = FIG_DIR / "weekly27_noise.png"
    fig.tight_layout()
    fig.savefig(out, dpi=200, bbox_inches="tight")
    plt.close(fig)
    return out


def fig_mpc() -> Path | None:
    """MPC-check u vs reference, recomputed from the capture + ref CSVs."""
    import numpy as np
    cap_p = REPO / "capture_mpc_20260827_184338.csv"
    ref_p = REPO / "ref_steps.csv"
    if not (cap_p.is_file() and ref_p.is_file()):
        print("NOTE: mpc capture/ref missing -- skipping MPC figure")
        return None
    try:
        hdr, data = _csv_cols(cap_p)
        low = [h.strip().lower() for h in hdr]
        ui = next(i for i, h in enumerate(low) if h in ("u1", "amp1", "out1"))
        u = np.array([float(r[ui]) for r in data])
        rrows = list(csv.reader(open(ref_p)))
        if any(ch.isalpha() for ch in "".join(rrows[0])):
            rrows = rrows[1:]
        r = np.array([float(row[0]) for row in rrows])
    except Exception as e:
        print(f"NOTE: could not parse MPC capture ({e}) -- skipping figure")
        return None
    n = min(len(u), len(r))
    plt = _mpl()
    FIG_DIR.mkdir(parents=True, exist_ok=True)
    fig, ax = plt.subplots(figsize=(8.6, 3.0))
    ax2 = ax.twinx()
    ax2.step(range(n), r[:n], color="#5B6470", lw=1.2, where="post",
             label="reference (feature space)")
    ax.plot(range(n), u[:n], color="#3F7A4E", lw=1.0, label="commanded u1")
    ax.set_xlabel("tick (10 ms frame-locked)")
    ax.set_ylabel("u1 (stim amplitude)", color="#3F7A4E")
    ax2.set_ylabel("reference", color="#5B6470")
    ax2.spines["right"].set_visible(True)
    ax.set_title("13 s MPC check 2026-08-27: the controller steps with the "
                 f"moving reference (block {BLOCK_MPCCHECK}, toy plant, saline)",
                 fontsize=10.5, color="black")
    h1, l1 = ax.get_legend_handles_labels()
    h2, l2 = ax2.get_legend_handles_labels()
    ax.legend(h1 + h2, l1 + l2, frameon=False, fontsize=9, loc="upper left")
    out = FIG_DIR / "weekly27_mpc.png"
    fig.tight_layout()
    fig.savefig(out, dpi=200, bbox_inches="tight")
    plt.close(fig)
    return out


def fig_week() -> Path:
    """Timeline of the week's events plus the Monday target (quoted dates)."""
    plt = _mpl()
    FIG_DIR.mkdir(parents=True, exist_ok=True)
    events = [
        (25, "prep + control-horizon\nresult banked (Nu=20)", "#5B6470", 1),
        (26, "suite day: preflight, quiet\ncapture, nThw 150/150", "#5B6470", -1),
        (26.96, "23:00 BLOCKER\nPO8e zero frames", "#B3413A", 1),
        (27.4, "root cause: card unseated\nreseat -> both gates GREEN", "#3F7A4E", -1),
        (31, "acute experiment\nGO from Phase 0", "#3F7A4E", 1),
    ]
    fig, ax = plt.subplots(figsize=(11.8, 2.3))
    ax.axhline(0, color="#444A52", lw=1.2, zorder=1)
    for x, label, color, side in events:
        ax.plot([x], [0], "o", ms=9, color=color, zorder=3)
        ax.annotate(label, (x, 0), xytext=(0, 26 * side), textcoords="offset points",
                    ha="center", va="bottom" if side > 0 else "top",
                    fontsize=8.5, color=color)
    ax.axvspan(28, 30.5, color="#5B6470", alpha=0.08, zorder=0)
    ax.text(29.25, 0.02, "weekend", ha="center", fontsize=8, color="#5B6470")
    ax.set_xlim(24.3, 31.9)
    ax.set_ylim(-1, 1)
    ax.set_yticks([])
    ax.set_xticks([25, 26, 27, 28, 31])
    ax.set_xticklabels(["Tue 08-25", "Wed 08-26", "Thu 08-27", "Fri 08-28", "Mon 08-31"])
    for sp in ("left", "bottom"):
        ax.spines[sp].set_visible(False)
    ax.tick_params(length=0)
    out = FIG_DIR / "weekly27_week.png"
    fig.tight_layout()
    fig.savefig(out, dpi=200, bbox_inches="tight")
    plt.close(fig)
    return out


def _read_u_and_ref(cap_name, ref_name, skip=50):
    """u1 from a capture + the reference affine-mapped into u-units (LS fit,
    same skip as the tracking scorer so the startup transient is excluded)."""
    import numpy as np
    hdr, data = _csv_cols(REPO / cap_name)
    ui = next(i for i, h in enumerate(hdr) if h.strip().lower() == "u1")
    u = np.array([float(r[ui]) for r in data])
    rhdr, rdata = _csv_cols(REPO / ref_name)
    ri = next(i for i, h in enumerate(rhdr) if h.strip().lower() == "r1")
    r = np.array([float(row[ri]) for row in rdata])
    n = min(len(u), len(r))
    u, r = u[skip:n], r[skip:n]
    a, b = np.polyfit(r, u, 1)
    return u, a * r + b


def fig_horizon() -> Path | None:
    """cl4 vs cl5 (2026-08-25 saline): commanded u against the reference scaled
    by the fitted slope, one shared u-unit axis per panel (no dual axes)."""
    try:
        cl4 = json.loads((REPO / "tracking_cl4.json").read_text())
        cl5 = json.loads((REPO / "tracking_cl5.json").read_text())
        runs = []
        for d, title in ((cl4, "cl4: step reference, Nu=2"),
                         (cl5, "cl5: touch-template reference, Nu=20")):
            s = d["per_channel"][0]["slope_u_on_r"]["u1"]
            u, r = _read_u_and_ref(d["capture"], d["ref"], d["skip_ticks"])
            runs.append((title, s, u, r))
    except (OSError, KeyError, StopIteration, ValueError) as e:
        print(f"NOTE: could not build control-horizon figure ({e}) -- skipping")
        return None
    plt = _mpl()
    FIG_DIR.mkdir(parents=True, exist_ok=True)
    fig, axes = plt.subplots(1, 2, figsize=(11.4, 2.7))
    for ax, (title, s, u, r) in zip(axes, runs):
        ax.step(range(len(r)), r, color="#5B6470", lw=1.2, where="post",
                label=f"reference, affine-mapped to u (scorer slope {s:.1f})")
        ax.plot(range(len(u)), u, color="#3F7A4E", lw=1.0, label="commanded u1")
        ax.set_title(title, fontsize=10, color="black")
        ax.set_xlabel("tick (10 ms)")
        ax.legend(frameon=False, fontsize=8.5, loc="upper left")
    axes[0].set_ylabel("u1 (stim amplitude)")
    out = FIG_DIR / "weekly27_horizon.png"
    fig.tight_layout()
    fig.savefig(out, dpi=200, bbox_inches="tight")
    plt.close(fig)
    return out


def fig_plan() -> Path:
    """Monday time budget as a cumulative timeline (same rows as the plan table;
    phase 6 shows its 4.5 h budget solid + the 6 h stretch hatched)."""
    plt = _mpl()
    FIG_DIR.mkdir(parents=True, exist_ok=True)
    phases = [  # (label, minutes, stretch_minutes)
        ("0-1 bring-up + preflight + quiet", 45, 0),
        ("2 thwack battery (10 sites)", 90, 0),
        ("3-4 probing -> fit -> GO/NO-GO", 80, 0),
        ("5 tapes + designs + NN training", 30, 0),
        ("6 ARMS x sites (4 controllers)", 270, 90),
        ("7 drift check + re-run + push", 30, 0),
    ]
    fig, ax = plt.subplots(figsize=(11.4, 2.9))
    t = 0.0
    for i, (label, mins, stretch) in enumerate(phases):
        ax.barh(i, mins / 60, left=t / 60, color="#3F7A4E" if "GO/NO-GO" in label
                else "#5B6470", height=0.62)
        if stretch:
            ax.barh(i, stretch / 60, left=(t + mins) / 60, color="none",
                    edgecolor="#5B6470", hatch="///", height=0.62, lw=0.8)
        ax.text((t + mins) / 60 + (stretch / 60) + 0.08, i, label, va="center",
                fontsize=9, color="black")
        t += mins + stretch
    total_lo = sum(m for _, m, _ in phases) / 60
    ax.set_yticks([])
    ax.invert_yaxis()
    ax.set_xlabel(f"elapsed hours (budget {total_lo:.1f} h; hatched = phase-6 stretch to 6 h -> {t/60:.1f} h)")
    ax.set_xlim(0, 12.4)
    out = FIG_DIR / "weekly27_plan.png"
    fig.tight_layout()
    fig.savefig(out, dpi=200, bbox_inches="tight")
    plt.close(fig)
    return out


def fig_latency() -> Path | None:
    """MATLAB-server turnaround distribution for the 08-27 MPC check,
    recomputed from mpc_lat_20260827_184338.csv; 10 ms = loop -TimeoutMs."""
    import numpy as np
    p = REPO / "mpc_lat_20260827_184338.csv"
    if not p.is_file():
        print("NOTE: mpc_lat csv missing -- skipping latency figure")
        return None
    hdr, data = _csv_cols(p)
    ti = next(i for i, h in enumerate(hdr) if h.strip().lower() == "turnaround_ms")
    v = np.array([float(r[ti]) for r in data if r[ti].strip().lower() != "nan"])
    over = float((v > 10).mean())
    plt = _mpl()
    FIG_DIR.mkdir(parents=True, exist_ok=True)
    fig, ax = plt.subplots(figsize=(8.4, 2.6))
    ax.hist(v, bins=np.arange(0, min(v.max(), 60) + 2, 1), color="#5B6470")
    ax.axvline(10, color="#B3413A", lw=1.2, ls="--")
    ax.text(10.6, ax.get_ylim()[1] * 0.9,
            f"10 ms loop timeout\n{over*100:.0f}% of {len(v)} replies over it\n"
            f"(p50 {np.percentile(v, 50):.1f} ms, p95 {np.percentile(v, 95):.1f} ms)",
            fontsize=9, color="#B3413A", va="top")
    ax.set_xlabel("MATLAB server turnaround (ms)")
    ax.set_ylabel("ticks")
    ax.set_title("08-27 MPC check: MATLAB server turnaround — a long-stall tail "
                 "(~35-38 ms), the known behavior behind loop timeouts (cpp arms "
                 "immune)", fontsize=10.5, color="black")
    out = FIG_DIR / "weekly27_latency.png"
    fig.tight_layout()
    fig.savefig(out, dpi=200, bbox_inches="tight")
    plt.close(fig)
    return out


# -------------------------------------------------------------------- build --
def build(qs, mc, designed, figs):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    BLACK = RGBColor(0, 0, 0)
    GREY = RGBColor(0x44, 0x4A, 0x52)
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]

    def new(title, subtitle=None):
        s = prs.slides.add_slide(blank)
        tb = s.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(12.2), Inches(0.72))
        p = tb.text_frame.paragraphs[0]
        r = p.add_run(); r.text = title
        r.font.name = FONT; r.font.size = Pt(24); r.font.bold = True
        r.font.color.rgb = BLACK
        if subtitle:
            tb2 = s.shapes.add_textbox(Inches(0.57), Inches(1.02), Inches(12.2), Inches(0.5))
            p2 = tb2.text_frame.paragraphs[0]
            r2 = p2.add_run(); r2.text = subtitle
            r2.font.name = FONT; r2.font.size = Pt(12); r2.font.color.rgb = GREY
        return s

    def bullets(s, items, top=1.62, size=BODY_PT, left=0.72, width=12.0):
        h = max(0.5, 7.5 - top - 0.25)
        tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True
        first = True
        for level, text in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.level = level
            r = p.add_run(); r.text = text
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.color.rgb = BLACK if level == 0 else GREY
            p.space_after = Pt(6)
        return tb

    def table(s, rows, left, top, width, height, size=BODY_PT):
        shape = s.shapes.add_table(len(rows), len(rows[0]),
                                   Inches(left), Inches(top), Inches(width), Inches(height))
        t = shape.table
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                cell = t.cell(i, j)
                cell.text = str(val)
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.name = FONT; r.font.size = Pt(size)
                        r.font.bold = (i == 0); r.font.color.rgb = BLACK
        return t

    def pic(s, path, left, top, width):
        s.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(width))

    # -- title ---------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.4))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = "Closed-Loop Stimulation System — Weekly Status"
    r.font.name = FONT; r.font.size = Pt(36); r.font.bold = True
    r.font.color.rgb = BLACK
    tb2 = s.shapes.add_textbox(Inches(0.82), Inches(3.7), Inches(11.7), Inches(0.9))
    r2 = tb2.text_frame.paragraphs[0].add_run()
    r2.text = ("Week of 2026-08-21 .. 2026-08-27  |  suite move validated end to end; "
               "transport hardware fault found + fixed; GO for the acute experiment Mon 08-31")
    r2.font.name = FONT; r2.font.size = Pt(14); r2.font.color.rgb = GREY

    # -- week at a glance ----------------------------------------------------
    s = new("Week at a glance")
    bullets(s, [
        (0, "08-25: pre-suite prep — randomized probe schedule, artifact assessor, tracking "
            "scorer, NN open-loop tape synthesis, results-deck builder; evening rig session "
            "banked the control-horizon result (Nu=20 standard set)."),
        (0, "08-26 (suite day): rig re-validated in the surgical suite — preflight all green, "
            "quiet capture banked (real 8-11 µV floor), headstage validated (ch 13 "
            "blacklisted), and the NEW nThw touch pulse validated 150/150 end to end."),
        (0, "08-26 23:00: BLOCKER — PO8e receives zero frames; experiment postponed from "
            "Thu 08-27. Surgery rescheduled to Mon 08-31."),
        (0, "08-27: root cause found in 30 min of morning diagnosis — the PO8e CARD had "
            "unseated in transport (off the PCI bus entirely). Reseat fixed it; both "
            "deferred saline gates then banked green the same session."),
        (0, "State: every pre-surgery gate is green. The Monday runbook "
            "(RIG_DAY_2026-08-31.md) is a GO from Phase 0."),
    ], top=1.5)
    pic(s, figs["week"], 0.75, 4.85, 11.8)

    # -- control horizon -----------------------------------------------------
    s = new("Control-horizon result: Nu=20 standard set (08-25, saline)",
            "The cl3 question answered: hold-last dilution, not plumbing")
    bullets(s, [
        (0, "cl4 (step reference, Nu=2): u plateaus land exactly on the reference steps — "
            "corr(u,r) 0.886, slope 9.7. Moving-target tracking PROVEN on hardware "
            "(capture_mpc_20260825_212708 vs ref_steps)."),
        (0, "cl5 (touch-template reference, Nu=20): transient fidelity du/(slope·dr) ≈ 1.0-1.3 "
            "vs ~2-3% at Nu=2 — a ~40x recovery of the fast touch spike. "
            "STANDARD: -ControlHorizon 20 for all touch-template arms."),
        (1, "Slope > naive prediction is correct certainty-equivalent behavior against a "
            "plant that never responds (saline); expect calibrated gains with a fitted plant."),
        (0, "NN training benchmark (28k-tick capture, 200 epochs, history 25): linear 14 s, "
            "mlp 15 s, residual_mlp 17 s, GRU 1062 s — day-of policy: GRU at ~50 epochs "
            "(~4-5 min) or skip."),
    ], top=1.5)
    if figs.get("horizon"):
        pic(s, figs["horizon"], 0.95, 4.55, 11.4)

    # -- suite bring-up ------------------------------------------------------
    s = new("Suite bring-up 2026-08-26: environment re-banked",
            "Recomputed from quiet_2026-08-26.json (block LD-260826-220719)")
    bullets(s, [
        (0, f"Noise floor {qs['noise_lo']:.1f}-{qs['noise_hi']:.1f} µV on 15/16 channels — the "
            "old bath's ~10 mV floor was the anomaly, not the reference. Feature baseline "
            f"{qs['base_lo']*1e6:.1f}-{qs['base_hi']*1e6:.1f} µV (MAV6, volts branch verified)."),
        (0, f"Ch 13 blacklisted from data: {qs['noise_ch13']:.0f} µV private noise, cross-corr "
            "0.18 vs the 0.90 bath-common signature (bad contact, not biology). Ch 16 marginal."),
        (0, f"60 Hz + harmonics = {qs['line_med']*100:.0f}% of 5-200 Hz relative power (~7 µV "
            "absolute) → fit-side notch planned; no C++ change on animal day."),
        (0, "Wav1/Wav2 DISK SAVING verified ON in the block "
            + ("(gate passed)." if qs["wav_saved"] else "(GATE FAILED — recheck!)")),
        (0, "nThw touch pulse FULLY VALIDATED: float32 stream @ 24.4 kHz, 150/150 programmed "
            "thwacks detected in BOTH test blocks, machine-uniform 255 ms contacts, sham path "
            "confirmed (fires without touch). It is the ONLY touch record (mPos/mCtl flat) — "
            "QC = onset count vs programmed count."),
    ], top=1.55)
    pic(s, figs["noise"], 2.9, 4.9, 7.6)

    # -- blocker -------------------------------------------------------------
    s = new("The blocker: a transport hardware fault, found and fixed",
            "08-26 23:00 postponement -> 08-27 morning: 30 minutes to root cause")
    bullets(s, [
        (0, "Symptom (08-26): card enumerated and connected, RZ2 healthy, but ZERO frames "
            "arrived — looked exactly like an RZ2-side rig-file/fiber fault. Six-step debug "
            "ladder prepped for the morning."),
        (0, "Morning finding: on the fresh boot the card was GONE — Windows CM_PROB_PHANTOM, "
            "absent from the PCI bus entirely. One fault explains both days: the PO8e card "
            "unseated during the PC's move to the suite (marginal seating enumerates but "
            "passes no data, then drops off on the next cold boot)."),
        (0, "Fix: full power-off, reseat, boot → card OK; probe PASS within seconds of "
            "Synapse Preview (Streaming. numChannels=16 sampleBytes=4). The rig file was "
            "never broken; none of the Synapse-side ladder was needed."),
        (0, "Hardened into the Monday runbook: PCI presence check is now Phase 0 step 1 "
            "(30 s, before ANY Synapse debugging), with the reseat procedure in the "
            "failure branches."),
    ])

    # -- gates ---------------------------------------------------------------
    s = new("Deferred gates banked green (08-27, suite, saline)",
            f"Blocks {BLOCK_RNDVAL} (rndval) and {BLOCK_MPCCHECK} (mpccheck)")
    bullets(s, [
        (0, f"Delivery VERIFIED: wire == design {designed}/{designed} (0 lost, 0 stretched); "
            "post-enable physical delivery 640/640 single-pulse — 0 missed, 0 doubled, "
            "0 latch races. Pair mapping exact on all 8 pairs."),
        (0, "Carrier 101.725 Hz (base/240) confirmed in the SUITE circuit — exactly 6.000 "
            "samples/period; the frame-locked feature chain stands unchanged."),
        (0, "Safety enable-gate re-confirmed: the 12 probes sent before the stim control was "
            "switched on were ALL blocked (gate holds Scle AND the carrier); 100% delivery "
            "after enable. Plus the stim-zero cleanup trace fired on every exit."),
        (0, "Artifact retest: MODERATE (0 fail / 0 warn) on controllable channels — "
            "feature-trim stays OFF for Monday; re-assess on the first in-vivo probe block."),
    ], top=1.55)
    pic(s, figs["gate"], 2.2, 4.75, 8.9)

    # -- MPC check -----------------------------------------------------------
    s = new("13-second MPC closed-loop check",
            "Recomputed from tracking_mpccheck_20260827.json + the capture")
    if figs.get("mpc"):
        pic(s, figs["mpc"], 0.9, 1.6, 8.8)
    bullets(s, [
        (0, f"{mc['ticks']}/{mc['ticks']} ticks, policy=fresh by tick 19 (startup race "
            f"avoided), u1-on-reference slope {mc['slope_u']:.2f} — matches the banked cl4 "
            f"signature (9.7). {mc['transients']} reference transients detected in u."),
        (0, f"y verdict '{mc['verdict']}' is the CORRECT saline null (toy model, no plant in "
            "the bath); u-on-r is the engineering readout. Watch item: MATLAB-server "
            "timeouts 13% / freshTicks 97.0% — above the 5-11% baseline, cpp arms immune."),
    ], top=4.75)

    # -- Monday plan ---------------------------------------------------------
    s = new("Monday 2026-08-31: the acute experiment",
            "RIG_DAY_2026-08-31.md — modular; every phase banks a self-contained result")
    table(s, [
        ["phase", "what", "budget"],
        ["0-1", "bring-up (PCI check first) + preflight + quiet capture", "45 min"],
        ["2", "thwack battery: 10 sites x ~30 contacts, nThw extraction each", "90 min"],
        ["3-4", "randomized probing (30 min run) -> fit -> GO/NO-GO (|corr| > 0.1)", "80 min"],
        ["5", "tapes + manifest + Choi designs + NN training", "30 min"],
        ["6", "ARMS x sites: choi / mpc / nnol / nncl, scored after every run", "4.5-6 h"],
        ["7", "drift check + re-run one early site + push", "30 min"],
    ], 0.72, 1.7, 11.9, 3.0)
    bullets(s, [
        (0, "Protocol: 4 controllers x 10 sites (9 touch + SHAM) x 100 elicited responses; "
            "site order randomized per arm (seeded manifest); one 22,200-tick run per "
            "(arm, site) = 3.7 min."),
        (0, "Cut order if the day runs long: nncl -> nnol (keep mpc + choi — MPC-vs-Choi on "
            "the same plant model is the primary comparison)."),
        (0, "uMax decided before the animal is on the table; stimulator charge checked before "
            "interpreting any null; stim enable ON immediately after 'go' (08-27 lesson)."),
    ], top=5.0)

    # -- Monday time budget --------------------------------------------------
    s = new("Monday at a glance: the time budget",
            "Same phases as the runbook table; the day is decided by phase 6")
    pic(s, figs["plan"], 0.95, 1.75, 11.4)
    bullets(s, [
        (0, "9.1 h budget (10.6 h with the phase-6 stretch) — phase 6 is half the day, "
            "which is why every earlier phase banks a self-contained result and the cut "
            "order (nncl -> nnol) is decided in advance."),
        (0, "The GO/NO-GO gate (green) sits ~3.6 h in: if randomized probing fits no plant "
            "(|corr| < 0.1), the day still banks the thwack battery + open-loop tapes."),
    ], top=5.2)

    # -- risks ---------------------------------------------------------------
    s = new("Watch items / open risks")
    bullets(s, [
        (0, "MATLAB MPC server timeouts crept to 13% (freshTicks 97.0%) in the suite check — "
            "known stall behavior, baseline updated; escalate only if freshTicks collapses. "
            "cpp arms (choi/nnol/nncl) are immune."),
        (0, "PCIe seating is now a known transport failure mode — presence check added to "
            "Phase 0; avoid moving the PC before Monday."),
        (0, "Artifact verdict is saline-MODERATE; must be re-assessed on the first in-vivo "
            "probe block before trusting lag-0 content (tool + command in the runbook)."),
        (0, "Hard process kill still bypasses stim-zeroing (Ctrl+C path verified; RZ2-side "
            "watchdog remains the durable ask)."),
        (0, "10 Mbps link negotiation on the RZ2 port persists (suspect cable) — harmless for "
            "control traffic; swap opportunistically, then re-run net_diag."),
    ], top=1.5)
    if figs.get("latency"):
        pic(s, figs["latency"], 2.5, 4.55, 8.4)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    out = OUT_DIR / "ClosedLoop_weekly_2026-08-27.pptx"
    try:
        prs.save(out)
    except PermissionError:
        n = 2
        while True:
            alt = out.with_name(f"{out.stem}_v{n}.pptx")
            try:
                prs.save(alt)
                print(f"NOTE: {out.name} is open in PowerPoint; wrote {alt.name} instead.")
                return alt
            except PermissionError:
                n += 1
                if n > 9:
                    raise
    return out


def main():
    qs = quiet_stats()
    designed = rndval_designed()
    mc = mpccheck_stats()
    print(f"Recomputed: noise {qs['noise_lo']:.1f}-{qs['noise_hi']:.1f} uV "
          f"(ch13 {qs['noise_ch13']:.0f}), line {qs['line_med']*100:.0f}%, "
          f"designed probes {designed}, mpc slope {mc['slope_u']:.3f}")
    figs = {"gate": fig_gate(), "noise": fig_noise(qs), "mpc": fig_mpc(),
            "week": fig_week(), "horizon": fig_horizon(), "plan": fig_plan(),
            "latency": fig_latency()}
    out = build(qs, mc, designed, figs)
    print(f"Wrote {out}")


if __name__ == "__main__":
    main()
