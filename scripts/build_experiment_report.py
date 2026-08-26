#!/usr/bin/env python
"""build_experiment_report.py -- next-day results deck for an acute experiment.

    python scripts/build_experiment_report.py --manifest experiment_manifest.json
        [--out ..\\PythonIntanAnalysis\\outputs\\Synthesis\\AcuteClosedLoop_<date>.pptx]
        [--charge-scale <uC per amp-unit per pulse>]

Everything on the slides is RECOMPUTED from run artifacts -- the manifest is
the only hand-written file. If an arm lists ref+capture but no tracking_json,
rig/tracking_metrics.py is run to produce it. Charge per arm is summed from
the capture's u columns (labelled a proxy unless --charge-scale is given).

Manifest schema (JSON):
{
 "experiment": "Acute closed-loop 2026-08-28", "date": "2026-08-28",
 "arms": [
   {"label": "rnd1", "kind": "probe", "capture": "capture_rig_runrnd1.csv",
    "probe_reports": ["probe_in1_report.json", "..."]},
   {"label": "choi1", "kind": "openloop", "capture": "capture_choi1.csv",
    "ref": "ref_touch.csv", "y_channels": [11], "lti": "plant_rnd1.lti",
    "tracking_json": "tracking_choi1.json", "tracking_png": "tracking_choi1.png"},
   {"label": "cl1", "kind": "closedloop", ...}
 ],
 "artifact_reports": ["artifact_rnd1_report.json"],
 "notes": ["free-text bullet", ...]
}

Formatting: Arial, black titles/headings (user rule: never blue). If the
output file is open in PowerPoint, a _vN copy is written instead.
"""

import argparse
import csv
import json
import os
import subprocess
import sys

import numpy as np

FONT = "Arial"
REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
VENV_PY = os.path.join(os.path.dirname(REPO),
                       "PythonIntanAnalysis", ".venv", "Scripts", "python.exe")


def load_capture_u(path):
    with open(path, newline="") as f:
        r = csv.reader(f)
        h = next(r)
        rows = np.array([[float(v) for v in row] for row in r])
    return rows[:, [i for i, c in enumerate(h) if c.startswith("u")]]


def charge_stats(capture, scale):
    U = load_capture_u(capture)
    total = float(U.sum())
    pulses = int((U > 1e-9).sum())
    per_ch = U.sum(axis=0)
    top = int(np.argmax(per_ch)) + 1
    if scale:
        return "%.3g uC total (%d pulse-ticks; most on u%d)" % (total * scale, pulses, top)
    return "%.4g amp-units*ticks (%d pulse-ticks; most on u%d)" % (total, pulses, top)


def ensure_tracking(arm):
    tj = arm.get("tracking_json")
    if tj and os.path.exists(tj):
        return tj
    if not (arm.get("ref") and arm.get("capture") and arm.get("y_channels")):
        return None
    tj = tj or ("tracking_%s.json" % arm["label"])
    png = arm.get("tracking_png") or ("tracking_%s.png" % arm["label"])
    cmd = [VENV_PY, os.path.join(REPO, "rig", "tracking_metrics.py"),
           "--ref", arm["ref"], "--capture", arm["capture"],
           "--y-channels"] + [str(c) for c in arm["y_channels"]] + \
          ["--label", arm["label"], "--json", tj, "--png", png]
    if arm.get("lti"):
        cmd += ["--lti", arm["lti"]]
    print("recomputing tracking for %s ..." % arm["label"])
    subprocess.run(cmd, cwd=REPO, check=False)
    arm["tracking_json"], arm["tracking_png"] = tj, png
    return tj if os.path.exists(tj) else None


def main():
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--manifest", required=True)
    ap.add_argument("--out", default=None)
    ap.add_argument("--charge-scale", type=float, default=None)
    args = ap.parse_args()

    os.chdir(REPO)
    man = json.load(open(args.manifest))
    date = man.get("date", "")
    out = args.out or os.path.join(os.path.dirname(REPO), "PythonIntanAnalysis",
                                   "outputs", "Synthesis",
                                   "AcuteClosedLoop_%s.pptx" % (date or "report"))

    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    BLACK = RGBColor(0, 0, 0)
    GREY = RGBColor(0x44, 0x4A, 0x52)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]

    def new(title, subtitle=None):
        s = prs.slides.add_slide(blank)
        box = s.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(12.2), Inches(0.72))
        p = box.text_frame.paragraphs[0]
        p.text = title
        p.font.size, p.font.bold, p.font.name, p.font.color.rgb = Pt(26), True, FONT, BLACK
        if subtitle:
            sb = s.shapes.add_textbox(Inches(0.57), Inches(1.02), Inches(12.2), Inches(0.5))
            sp = sb.text_frame.paragraphs[0]
            sp.text = subtitle
            sp.font.size, sp.font.name, sp.font.color.rgb = Pt(13), FONT, GREY
        return s

    def bullets(s, items, top=1.62, size=15.5, left=0.72, width=12.0):
        tf = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width),
                                  Inches(max(0.5, 7.5 - top - 0.25))).text_frame
        tf.word_wrap = True
        first = True
        for lvl, text in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = text
            p.level = lvl
            p.font.size = Pt(size if lvl == 0 else size - 1.5)
            p.font.name = FONT
            p.font.color.rgb = BLACK if lvl == 0 else GREY
            p.space_after = Pt(7)

    def table(s, rows, left, top, width, height, size=12):
        shp = s.shapes.add_table(len(rows), len(rows[0]), Inches(left), Inches(top),
                                 Inches(width), Inches(height)).table
        for r, row in enumerate(rows):
            for c, val in enumerate(row):
                cell = shp.cell(r, c)
                cell.text = str(val)
                para = cell.text_frame.paragraphs[0]
                para.font.size, para.font.name = Pt(size), FONT
                para.font.bold = (r == 0)
                para.font.color.rgb = BLACK

    # ---- title -------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    b = s.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.4))
    p = b.text_frame.paragraphs[0]
    p.text = man.get("experiment", "Acute closed-loop experiment")
    p.font.size, p.font.bold, p.font.name, p.font.color.rgb = Pt(38), True, FONT, BLACK
    b2 = s.shapes.add_textbox(Inches(0.95), Inches(3.6), Inches(11.5), Inches(1.2))
    tf = b2.text_frame
    tf.word_wrap = True
    for i, (txt, sz) in enumerate([
            ("Results: probing, model, open- vs closed-loop arms", 19),
            ("TDTClosedLoop  |  %s" % date, 13)]):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = txt
        para.font.size, para.font.name, para.font.color.rgb = Pt(sz), FONT, GREY

    # ---- protocol table ----------------------------------------------------
    s = new("Protocol", "arms in execution order; every number recomputed from run artifacts")
    rows = [["arm", "kind", "capture", "reference", "charge delivered"]]
    for arm in man.get("arms", []):
        rows.append([arm["label"], arm.get("kind", ""),
                     os.path.basename(arm.get("capture", "")) or "-",
                     os.path.basename(arm.get("ref", "")) or "-",
                     charge_stats(arm["capture"], args.charge_scale)
                     if arm.get("capture") and os.path.exists(arm["capture"]) else "-"])
    table(s, rows, 0.6, 1.7, 12.1, 0.4 * len(rows))
    if args.charge_scale is None:
        bullets(s, [(1, "charge shown as amp-units*ticks (no --charge-scale given); "
                        "multiply by uC-per-amp-unit-per-pulse for physical charge")],
                top=1.85 + 0.42 * len(rows), size=12)

    # ---- probe SNR ---------------------------------------------------------
    probe_rows = [["input (pair)", "best ch", "gain", "delay (ticks)",
                   "SNR vs p99 floor", "linearity R2", "verdict"]]
    for arm in man.get("arms", []):
        for rp in arm.get("probe_reports", []):
            if not os.path.exists(rp):
                probe_rows.append([rp, "MISSING", "", "", "", "", ""])
                continue
            d = json.load(open(rp))
            chans = d.get("channels", [])
            best = chans[0] if chans else {}
            snr = best.get("snr_vs_p99floor", float("nan"))
            lin = best.get("linearity_r2")
            probe_rows.append([
                d.get("input_1based", "?"), best.get("channel", "?"),
                "%.4g" % best.get("gain", float("nan")),
                best.get("delay_ticks", "?"),
                "%.2f" % snr,
                ("%.2f" % lin) if lin is not None else "-",
                "RESPONDING" if snr >= 3 else
                ("marginal" if snr >= 1.5 else "NO RESPONSE")])
    if len(probe_rows) > 1:
        s = new("Single-pulse probing", "randomized balanced-deck schedule; "
                "fit_impulse_model per input")
        table(s, probe_rows, 0.6, 1.7, 12.1, 0.38 * len(probe_rows))

    # ---- per-arm tracking slides ------------------------------------------
    for arm in man.get("arms", []):
        if arm.get("kind") not in ("openloop", "closedloop", "nn_openloop",
                                   "nn_closedloop"):
            continue
        tj = ensure_tracking(arm)
        s = new("Arm %s (%s)" % (arm["label"], arm.get("kind", "")),
                "reference vs achieved -- rig/tracking_metrics.py")
        if tj:
            d = json.load(open(tj))
            rows = [["ref col -> y ch", "NRMSE (ptp)", "r (best lag)", "lag",
                     "slope y/r", "track idx", "verdict"]]
            for ch in d["per_channel"]:
                rows.append(["r%d -> y%d" % (ch["ref_col"], ch["y_channel"]),
                             "%.1f%%" % (100 * ch["nrmse_ptp"]),
                             "%.3f" % ch["pearson_best"],
                             "%+d" % ch["best_lag_ticks"],
                             "%.3f" % ch["slope_y_on_r"],
                             "%.2f" % ch["tracking_index"]
                             if ch["tracking_index"] == ch["tracking_index"] else "-",
                             ch["verdict"]])
            rows.append(["OVERALL", "", "", "", "", "", d["overall"]])
            table(s, rows, 0.6, 1.62, 6.4, 0.38 * len(rows))
        else:
            bullets(s, [(0, "tracking metrics unavailable (no ref/capture/"
                            "y_channels in the manifest)")])
        png = arm.get("tracking_png")
        if png and os.path.exists(png):
            s.shapes.add_picture(png, Inches(7.2), Inches(1.55), width=Inches(5.7))

    # ---- artifact assessment ----------------------------------------------
    arts = [a for a in man.get("artifact_reports", []) if os.path.exists(a)]
    if arts:
        s = new("Stimulus artifact assessment", "rig/assess_artifact.py; verdict "
                "scores controllable (off-pair) channels only")
        rows = [["report", "store", "events", "worst off-pair std ratio", "verdict"]]
        for a in arts:
            d = json.load(open(a))
            wo = d.get("worst_off_pair") or {}
            rows.append([os.path.basename(a).replace("_report.json", ""),
                         d.get("store", "?"), d.get("n_events", "?"),
                         "%.1f (word %s ch %s)" % (wo.get("std_ratio", float("nan")),
                                                   wo.get("word", "?"),
                                                   wo.get("channel", "?"))
                         if wo else "-",
                         d.get("verdict", "?")])
        table(s, rows, 0.6, 1.7, 12.1, 0.4 * len(rows))
        for i, a in enumerate(arts[:2]):
            png = a.replace("_report.json", ".png")
            if os.path.exists(png):
                s.shapes.add_picture(png, Inches(0.7 + 6.3 * i),
                                     Inches(2.0 + 0.42 * len(rows)), width=Inches(5.9))

    # ---- notes / next steps ------------------------------------------------
    items = [(0, n) for n in man.get("notes", [])]
    for arm in man.get("arms", []):
        tj = arm.get("tracking_json")
        if tj and os.path.exists(tj):
            d = json.load(open(tj))
            items.append((1, "%s: %s" % (arm["label"], d["overall"])))
    if items:
        bullets(new("Notes and verdicts"), items)

    try:
        prs.save(out)
    except PermissionError:
        base, ext = os.path.splitext(out)
        n = 2
        while True:
            alt = "%s_v%d%s" % (base, n, ext)
            try:
                prs.save(alt)
                out = alt
                break
            except PermissionError:
                n += 1
    print("Wrote %s (%d slides)" % (out, len(prs.slides._sldIdLst)))
    return 0


if __name__ == "__main__":
    sys.exit(main())
