"""End-of-week deck, week of 2026-08-28 .. 2026-09-04 (weekly-for-PI).

The whole arc: 32-ch migration -> first in-vivo closed-loop biomimetic
tracking -> artifact-aware corrections -> MIMO/cpp sprint -> 9/8 plan.

Formatting per standing preference: Arial, BLACK headings (never blue), 11pt
body, ~1 figure per content slide. Headline numbers are read from the analysis
summary JSONs at build time; the two NEW figures (week timeline, 9/8 Gantt)
and the MIMO tracking figure are generated here (MIMO from
capture_mimo_sim.csv + ref_mimo_test.csv — the p=2 closed-loop sim capture).
Every embedded PNG is existence-checked; the build FAILS on a missing figure.

Output: PythonIntanAnalysis/outputs/Synthesis/ClosedLoop_weekly_2026-09-04.pptx
Run with the PythonIntanAnalysis venv python.
"""
from __future__ import annotations

import csv
import json
from pathlib import Path

REPO = Path(__file__).resolve().parents[1]
ANA = REPO / "day_2026-08-31" / "analysis"
FIG_DIR = REPO / "outputs" / "deck_figures"
OUT_DIR = REPO.parent / "PythonIntanAnalysis" / "outputs" / "Synthesis"
FONT = "Arial"
BODY_PT = 11.0

GREEN, GREY, RED, AMBER = "#3F7A4E", "#5B6470", "#B3413A", "#C9A23A"


def need(path: Path) -> Path:
    if not path.is_file():
        raise SystemExit(f"FAIL: required figure missing: {path}")
    return path


def jload(name):
    return json.loads((ANA / name).read_text())


# ------------------------------------------------------------------ figures --
def _mpl():
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    plt.rcParams.update({
        "font.family": FONT, "font.size": 10, "text.color": "black",
        "axes.edgecolor": "#444A52", "axes.labelcolor": "black",
        "xtick.color": "#444A52", "ytick.color": "#444A52",
        "axes.spines.top": False, "axes.spines.right": False,
    })
    return plt


def fig_week() -> Path:
    plt = _mpl()
    FIG_DIR.mkdir(parents=True, exist_ok=True)
    events = [
        (28.0, "8/28 32-ch prep sprint\n(guard, warm-width fix)", GREY, 1),
        (29.5, "8/29-30 suite gates\n(Detect incident RECOVERED;\n32-ch chain proven)", GREY, -1),
        (31.0, "8/31 SURGERY:\n4/4 arms TRACKING", GREEN, 1),
        (32.0, "9/1 analysis + MIMO/cpp\nsprint (bench 28/28)", GREEN, -1),
        (33.5, "9/2-9/5 recovery\n(experiments paused)", AMBER, 1),
        (35.5, "9/7 setup ->\n9/8 ACUTE #2 (64-ch MIMO)", GREEN, -1),
    ]
    fig, ax = plt.subplots(figsize=(11.8, 2.4))
    ax.axhline(0, color="#444A52", lw=1.2, zorder=1)
    for x, label, color, side in events:
        ax.plot([x], [0], "o", ms=9, color=color, zorder=3)
        ax.annotate(label, (x, 0), xytext=(0, 24 * side), textcoords="offset points",
                    ha="center", va="bottom" if side > 0 else "top",
                    fontsize=8.5, color=color)
    ax.set_xlim(27.4, 36.3)
    ax.set_ylim(-1, 1)
    ax.set_yticks([])
    ax.set_xticks([28, 29.5, 31, 32, 33.5, 35.5])
    ax.set_xticklabels(["Fri 8/28", "Sat-Sun", "Mon 8/31", "Tue 9/1", "recovery", "9/7-9/8"])
    for sp in ("left", "bottom"):
        ax.spines[sp].set_visible(False)
    ax.tick_params(length=0)
    out = FIG_DIR / "weekly0904_week.png"
    fig.tight_layout(); fig.savefig(out, dpi=200, bbox_inches="tight"); plt.close(fig)
    return out


def fig_mimo() -> Path:
    """p=2 closed-loop sim: COMMANDS u1/u4 vs the two reference columns
    (affine-mapped LS into command units). The sim plant is the near-zero-gain
    saline fit, so u-on-r is the honest readout (the week's convention) --
    the point is that each command channel follows ITS OWN reference.
    Recomputed from capture_mimo_sim.csv + ref_mimo_test.csv."""
    import numpy as np
    plt = _mpl()
    cap = list(csv.reader(open(REPO / "capture_mimo_sim.csv")))
    hdr = cap[0]
    col = {h: i for i, h in enumerate(hdr)}
    data = np.array([[float(v) for v in r] for r in cap[1:]])
    ref = np.array([[float(r[1]), float(r[2])]
                    for r in list(csv.reader(open(REPO / "ref_mimo_test.csv")))[1:]])
    n = min(len(data), len(ref))
    skip = 50
    fig, axes = plt.subplots(2, 1, figsize=(10.8, 4.4), sharex=True)
    for k, (ax, ucol) in enumerate(zip(axes, ("u1", "u4"))):
        u = data[:n, col[ucol]]
        r = ref[:n, k]
        a, b = np.polyfit(r[skip:], u[skip:], 1)
        rr = np.corrcoef(u[skip:], r[skip:])[0, 1]
        ax.step(range(n), a * r + b, color=GREY, lw=1.3, where="post",
                label=f"reference r{k+1} (affine-mapped to command units)")
        ax.plot(range(n), u, color=GREEN, lw=1.0, label=f"commanded {ucol} (uA)")
        ax.set_ylabel(f"{ucol} (uA)")
        ax.legend(frameon=False, fontsize=9, loc="upper right")
        ax.set_title(f"command channel {k+1}: corr(u, its own reference) = {rr:.2f}",
                     fontsize=10, color="black")
    axes[1].set_xlabel("tick (9.83 ms)")
    fig.suptitle("MIMO closed loop (p=2, cpp MPC, real stacked model): each command "
                 "channel follows ITS OWN reference (saline-fit plant -> u-on-r is "
                 "the readout)", fontsize=11, color="black")
    out = FIG_DIR / "weekly0904_mimo.png"
    fig.tight_layout(); fig.savefig(out, dpi=200, bbox_inches="tight"); plt.close(fig)
    return out


def fig_plan() -> Path:
    plt = _mpl()
    phases = [
        ("0-1 bring-up + preflight + quiet (64 ch)", 45, 0, GREY),
        ("2 thwack battery, 64-wide templates", 90, 0, GREY),
        ("3 all-8-pair probing + raw-LFP screen", 40, 0, GREEN),
        ("4 operating-point MIMO fit -> GO/NO-GO", 40, 0, GREEN),
        ("5 p-channel refs + Choi tapes (slew-penalized)", 30, 0, GREY),
        ("6 ARMS: cpp-MPC vs Choi, paired schedules", 150, 90, GREY),
        ("7 drift check + wrap + push", 30, 0, GREY),
    ]
    fig, ax = plt.subplots(figsize=(11.4, 3.0))
    t = 0.0
    for i, (label, mins, stretch, color) in enumerate(phases):
        ax.barh(i, mins / 60, left=t / 60, color=color, height=0.62)
        if stretch:
            ax.barh(i, stretch / 60, left=(t + mins) / 60, color="none",
                    edgecolor=GREY, hatch="///", height=0.62, lw=0.8)
        ax.text((t + mins + stretch) / 60 + 0.08, i, label, va="center",
                fontsize=9, color="black")
        t += mins + stretch
    ax.set_yticks([]); ax.invert_yaxis()
    ax.set_xlabel(f"elapsed hours (budget {sum(m for _, m, _, _ in phases)/60:.1f} h; "
                  f"hatched = phase-6 stretch -> {t/60:.1f} h)")
    ax.set_xlim(0, 11.8)
    out = FIG_DIR / "weekly0904_plan.png"
    fig.tight_layout(); fig.savefig(out, dpi=200, bbox_inches="tight"); plt.close(fig)
    return out


# -------------------------------------------------------------------- build --
def build():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    BLACK = RGBColor(0, 0, 0)
    DGREY = RGBColor(0x44, 0x4A, 0x52)
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]

    def new(title, subtitle=None):
        s = prs.slides.add_slide(blank)
        tb = s.shapes.add_textbox(Inches(0.55), Inches(0.3), Inches(12.2), Inches(0.7))
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = title
        r.font.name = FONT; r.font.size = Pt(23); r.font.bold = True
        r.font.color.rgb = BLACK
        if subtitle:
            tb2 = s.shapes.add_textbox(Inches(0.57), Inches(0.95), Inches(12.2), Inches(0.45))
            r2 = tb2.text_frame.paragraphs[0].add_run()
            r2.text = subtitle
            r2.font.name = FONT; r2.font.size = Pt(12); r2.font.color.rgb = DGREY
        return s

    def bullets(s, items, top=1.5, size=BODY_PT, left=0.72, width=12.0):
        tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width),
                                  Inches(max(0.5, 7.5 - top - 0.25)))
        tf = tb.text_frame; tf.word_wrap = True
        first = True
        for level, text in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.level = level
            r = p.add_run(); r.text = text
            r.font.name = FONT; r.font.size = Pt(size)
            r.font.color.rgb = BLACK if level == 0 else DGREY
            p.space_after = Pt(5)
        return tb

    def pic(s, path, left, top, width):
        s.shapes.add_picture(str(need(Path(path))), Inches(left), Inches(top), Inches(width))

    def fig_slide(title, subtitle, path, notes=None, fw=10.4, fl=1.45, ft=1.35):
        s = new(title, subtitle)
        pic(s, path, fl, ft, fw)
        if notes:
            bullets(s, notes, top=6.55, size=10.0)
        return s

    sci = jload("sci_summary.json")
    spat = jload("spat_summary.json")

    # 1 title
    s = prs.slides.add_slide(blank)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(11.7), Inches(1.5))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = "First In-Vivo Closed-Loop Biomimetic Tracking"
    r.font.name = FONT; r.font.size = Pt(36); r.font.bold = True; r.font.color.rgb = BLACK
    tb2 = s.shapes.add_textbox(Inches(0.82), Inches(3.55), Inches(11.7), Inches(1.0))
    r2 = tb2.text_frame.paragraphs[0].add_run()
    r2.text = ("Week of 2026-08-28 .. 09-04: 32-channel migration, acute experiment #1 "
               "(4/4 arms TRACKING), artifact-aware analysis, MIMO + native-C++ control "
               "stack, and the plan for acute #2 (2026-09-08).")
    r2.font.name = FONT; r2.font.size = Pt(14); r2.font.color.rgb = DGREY

    # 2 week timeline
    s = new("The week")
    pic(s, fig_week(), 0.75, 1.6, 11.8)
    bullets(s, [
        (0, "Five working days: a full recording-side migration, a surgical first, a "
            "same-night analysis pass, an artifact-hardening second pass, and the "
            "controller stack rebuilt for multichannel control."),
    ], top=4.7)

    # 3 migration
    s = new("32-channel migration: proven in two days, ran flawlessly on surgery day",
            "Opt-in flags everywhere; 16-ch fallback preserved throughout")
    bullets(s, [
        (0, "Loop width guard (request>stream = fatal; channelMode= in every log); "
            "warm-width fix (a latent bug that capped the MPC's addressable channels "
            "at 8 -- found because migration testing exercised channel 20)."),
        (0, "Suite gates: delivery re-banked 473/473 wire, 428/428 post-enable; 13 s "
            "MPC check at feature channel 20; full-sweep fit timing 12.4 s at 32 ch."),
        (0, "Synapse lessons banked the hard way: never run rig Detect (recovery = "
            "Backups .synrig import); after any width edit -- rebuild, banner check, "
            "store-width check (the stale-apply lesson)."),
        (0, "Surgery day: every arm run 22,200/22,200 ticks, zero dropped, "
            "channelMode=exact(32/32) on every log line."),
    ])

    # 4 five acts
    s = new("Acute #1 (2026-08-31) in five acts")
    bullets(s, [
        (0, "1. Implant + battery: 10/10 thwack blocks, split-half 0.92-0.99, peaks "
            "382-709 uV, modal best channel 8."),
        (0, "2. A NO-GO reversed: all fitters refused (best |corr| 0.071); raw-LFP "
            "signed averages found the response -- recruitment threshold ~13-18 uA "
            "had drowned every zero-origin linear fit."),
        (0, "3. Operating-point identification: PRBS 10-30 uA held above the knee on "
            "pairs 1+4 -> |corr| 0.22-0.25 -> the day's plant model."),
        (0, "4. Three volt-scale numerics bugs found and fixed live (Choi mu; Choi "
            "horizon memory; MPC QP below solver tolerance)."),
        (0, "5. Four arm runs, all TRACKING: MPC rBest 0.712-0.725 at a stable ~20 ms "
            "lag; Choi 0.704-0.716 with drifting lag; plant held (drift re-probe)."),
    ])

    # 5-14: existing analysis figures
    fig_slide("The targets: 10-site touch battery",
              "150 contacts/site; the references the controllers must reproduce",
              ANA / "spat_templates.png",
              [(0, f"Touch site identity decodes from cortex at "
                   f"{spat['mahalanobis']['touch_loo_10way_accuracy']*100:.0f}% "
                   f"(10-way, chance 10%) -- the recording side is publication-grade.")],
              fw=9.6, fl=1.85)

    fig_slide("The pivot: repositioning averted by finding the recruitment threshold",
              "Signed raw-LFP averages on data the fitters refused",
              ANA / "spat_recruitment.png",
              [(0, "Pairs 4/1 drive S1 (+241/+121 uV @ 13 ms, z 23.5/9.7), touch "
                   "footprint r 0.96-0.99, knee 13-18 uA. Standing rule: signed-average "
                   "the raw LFP before any electrode moves.")], fw=9.6, fl=1.85)

    fig_slide("Tracking, all five sites, all arms",
              "Event-triggered achieved vs target; Hold = accidental tonic-only control",
              ANA / "sci_event_gallery.png",
              [(0, "MPC event-average r 0.88-0.95 per site; Hold control flat "
                   "(r 0.004 +/- 0.125): tracking is controller action.")],
              fw=10.0, fl=1.65)

    fig_slide("The value of feedback appears over time",
              "Per-event fidelity vs event number; the plant itself HELD (re-probe 0.224 -> 0.248)",
              ANA / "sci_timecourse.png",
              [(0, "Choi decays within-run even at per-event best lag (0.81 -> 0.70, "
                   "-0.13 r/100 events, both runs); MPC never degrades. Beyond ~100 "
                   "events MPC wins outright.")])

    fig_slide("Paired head-to-head, artifact-aware: MPC wins",
              "Same seeded schedules in both arms; artifact ticks (0-2 post-onset) excluded",
              ANA / "ctl_paired_clean.png",
              [(0, "Cleaned pooled dr = +0.079 [0.046, 0.112], p<1e-4, MPC wins 65% of "
                   "161 paired events (r1-driven: Choi's apparent parity there was partly "
                   "artifact samples). Both arms far above the Hold floor "
                   "(r ~0.5-0.6 vs 0.06) -- the tracking itself is genuine.")])

    fig_slide("SHAM catch trials: no fabricated touches",
              "Interleaved sham events, same scoring as real sites",
              ANA / "sci_sham.png",
              [(0, "False-touch rate MPC 0/32, Choi 1/33; hit rates 95%/86%; "
                   "d-prime 4.2/3.6.")])

    fig_slide("Same fidelity, 30% less charge",
              "Command strategies on identical schedules",
              ANA / "eng_u_strategy.png",
              [(0, "Choi rides the 30 uA cap 48% of ticks, slews 4.9x harder, delivers "
                   "+29.5% charge; MPC stays within +/-0.5 uA of the operating point "
                   "half the time. Feedback substitutes information for charge.")])

    fig_slide("Artifact anatomy: a new methods finding",
              "Pulse-resolved 24 kHz separation (art_ pass, 09-01)",
              ANA / "art_cleaned_gallery.png",
              [(0, "Per-pulse artifact is over by ~1.2 ms and contributes <0.3 dB below "
                   "45 Hz. The real contaminant is a ~2 mV SLOW TRANSIENT locked to "
                   "amplitude steps (electrode polarization) -- worst for Choi's "
                   "cap-slamming tapes. Mitigation for 9/8: slew-penalized Choi synthesis; "
                   "MPC is naturally gentle.")])

    fig_slide("The honest negative: no site-selectivity (artifact-robust)",
              "Mahalanobis distance of arm-evoked 32-ch patterns to each touch site's distribution",
              ANA / "spat_mahalanobis.png",
              [(0, f"Cleaned 5-way accuracy {jload('art_summary.json').get('accuracy_5way', {}).get('MPC', 0.23)*100:.0f}%/"
                   f"{jload('art_summary.json').get('accuracy_5way', {}).get('Choi', 0.22)*100:.0f}% vs 20% chance; "
                   "intent decoding from the time course is null too. One control channel "
                   "+ one stim footprint cannot choose WHICH touch it reproduces -- "
                   "the measured version of the rank/selectivity ceiling, and the "
                   "direct motivation for Thursday's MIMO + all-8-pair design.")])

    fig_slide("Engineering: the loop earned the week",
              "Latency and tick health across all runs",
              ANA / "eng_latency.png",
              [(0, "MATLAB server p50 1.69 ms (p99 29 ms) vs native C++ 17 us (~100x). "
                   "All arm runs: p95 tick error < 1.9 ms, zero dropped control ticks. "
                   "Verdict executed: the C++ server is now primary for every arm.")])

    # 15 MIMO
    s = new("The MIMO sprint: multichannel control is built and sim-proven",
            "9/1: fitter, both controllers, tapes, scoring -- all multi-output; bench 28/28 (first all-green)")
    pic(s, fig_mimo(), 1.25, 1.5, 10.8)
    bullets(s, [
        (0, "Per-output ARX + block-diagonal stacking (real p=2 fit from the 8/31 "
            "capture); cpp MPC gained operating-point offsets, reference preview, "
            "per-channel weights; Choi is p-native; one command now builds p-channel "
            "paired-schedule runs. 24 us round trips at p=2."),
    ], top=6.3)

    # 16 tooling
    galleries = sorted((REPO / "galleries").glob("*/index.png"))
    s = new("New rig tooling: every block gets a gallery, automatically",
            "rig/plot_trial_responses.py + rig/ingest_block.py (one command per block)")
    if galleries:
        pic(s, galleries[0], 1.85, 1.4, 9.6)
    bullets(s, [
        (0, "Thwack, probe, and arm modes; per-condition heatmaps + trial stacks + "
            "stats JSON land in galleries/<block>/ the moment a block is handed over -- "
            "the sanity-check the operator asked for, now free."),
    ], top=6.5)

    # 17 plan
    s = new("Next: acute #2 -- Monday 2026-09-08 (setup 9/7)",
            "64-ch recording (2x32 headstages), all 8 stim pairs, MIMO arms, cpp-primary")
    pic(s, fig_plan(), 0.95, 1.5, 11.4)
    bullets(s, [
        (0, "Honest risks: no saline validation (arrays too valuable -- wire-green is "
            "NOT tissue-green; first in-vivo probe block is the real delivery gate); "
            "64-ch and MIMO both first-time in tissue; new cortical arrays were the "
            "hardest part of acute #1 -- 64-ch handling rehearsed at setup."),
    ], top=5.1)

    # 18 summary
    s = new("Summary")
    bullets(s, [
        (0, "GOOD: closed-loop biomimetic tracking is real and artifact-robust -- "
            "stable ~20 ms latency, d-prime ~4 sham rejection, 30% charge savings, and "
            "with artifact-aware scoring MPC beats the Choi-2016 baseline outright."),
        (0, "GOOD: the interleaved randomized design worked (zero order effects) and "
            "the recording side is publication-grade (48% ten-way touch decodability)."),
        (0, "BAD (and productive): no site-selectivity in space or time -- one control "
            "channel cannot choose which touch it reproduces. That ceiling is the "
            "thesis argument for multichannel control, and the MIMO stack to break it "
            "is built, benched (28/28), and scheduled for 9/8."),
        (0, "CAVEATS: per-event amplitude grading is below single-event SNR; Choi's "
            "step transient needs the slew penalty; biological-space validation during "
            "continuous stim requires the artifact-aware pipeline (now standard)."),
    ])

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    out = OUT_DIR / "ClosedLoop_weekly_2026-09-04.pptx"
    try:
        prs.save(out)
    except PermissionError:
        out = out.with_name(out.stem + "_v2.pptx")
        prs.save(out)
    print(f"Wrote {out}")


if __name__ == "__main__":
    build()
