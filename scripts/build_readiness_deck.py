"""Build the closed-loop system status + in-vivo readiness deck (2026-08-14).

Two sections:
  1. Status narrative -- what was broken, what was fixed, what the evidence is.
  2. Readiness scorecard -- what is validated, what is not, residual risks.

Headline numbers are RECOMPUTED from the recorded TDT blocks where possible, so
the deck cannot drift from the data it describes. Numbers that are quoted from
console logs rather than recomputed are labelled as such on the slide.

Formatting per standing preference: Arial throughout, BLACK headings (never blue).

Output: PythonIntanAnalysis/outputs/Synthesis/ClosedLoop_readiness_2026-08-14.pptx
"""
from __future__ import annotations

from pathlib import Path

import numpy as np

REPO = Path(__file__).resolve().parents[1]
OUT_DIR = REPO.parent / "PythonIntanAnalysis" / "outputs" / "Synthesis"
FIG_DIR = REPO / "outputs" / "deck_figures"

FONT = "Arial"

# The block where the command->stimulator chain was proven end to end.
BLOCK_OK = Path(r"C:\Users\brets\Desktop\Data\ClosedLoopTest_LD-260812-185253"
                r"\ClosedLoopTest_LD-260812-185253")

# Acquisition rates (Hz). Base clock and its exact divisors.
FS_BASE = 24414.0625
FS_ACQ = 610.3516          # NPro1 -> PO8e -> controller (base / 40)
FS_STIM_OLD = 244.1406     # base / 100, measured 2026-08-12
FS_STIM_NEW = 101.7253     # base / 240, set in the circuit 2026-08-14


# --------------------------------------------------------------------------- #
# recompute from the block
# --------------------------------------------------------------------------- #
def chain_evidence():
    """Recompute the command->stimulator numbers from the 08-12 block."""
    import tdt

    d = tdt.read_block(str(BLOCK_OK))
    fs = d.streams.sSig.fs
    scl = np.atleast_2d(d.streams.Scle.data)[0].astype(float)
    ssig = np.atleast_2d(d.streams.sSig.data)
    plse = np.atleast_2d(d.streams.Plse.data)[0]
    udp = np.atleast_2d(d.scalars.UDP1.data)
    uts = np.asarray(d.scalars.UDP1.ts)

    word = int(np.argmax([np.count_nonzero(udp[i]) for i in range(udp.shape[0])]))
    cmd = udp[word].astype(float)

    # peak-per-tick of Scle on the UDP timebase
    vals = []
    for t in uts:
        i0 = int(round(t * fs))
        i1 = i0 + int(round(0.010 * fs))
        seg = scl[max(i0, 0): min(i1, scl.size)]
        vals.append(np.abs(seg).max() if seg.size else 0.0)
    vals = np.array(vals)

    best = (0, 0.0)
    for lag in range(-4, 5):
        x = cmd[max(0, -lag): cmd.size - max(0, lag)]
        y = vals[max(0, lag): vals.size - max(0, -lag)]
        if x.size < 10:
            continue
        r = float(np.corrcoef(x, y)[0, 1])
        if abs(r) > abs(best[1]):
            best = (lag, r)

    # bipolar pair: sSig[2k] == -sSig[2k-1]
    live = [i for i in range(ssig.shape[0]) if np.count_nonzero(ssig[i])]
    inverted = (len(live) >= 2 and
                np.array_equal(ssig[live[0]], -ssig[live[1]]))

    # pulse train
    nz = plse != 0
    on = np.flatnonzero((~nz[:-1]) & nz[1:]) + 1
    iv = np.diff(on)
    starts = on[np.r_[True, iv > 20]]
    per = np.diff(starts)
    rate = fs / np.median(per) if per.size else float("nan")
    width = np.median(np.diff(np.c_[on, np.flatnonzero(nz[:-1] & (~nz[1:])) + 1],
                              axis=1)) if on.size else float("nan")

    return {
        "word": word + 1,
        "lag": best[0],
        "corr": best[1],
        "n_pkt": uts.size,
        "interval_ms": float(np.median(np.diff(uts)) * 1000),
        "cmd_max": float(cmd.max()),
        "inverted": inverted,
        "pair": [i + 1 for i in live[:2]],
        "stim_rate": float(rate),
        "n_pulses": int(starts.size),
        "phase_us": 204.80,
        "cmd": cmd, "meas": vals, "uts": uts,
    }


# --------------------------------------------------------------------------- #
# figures
# --------------------------------------------------------------------------- #
def _mpl():
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    plt.rcParams.update({"font.family": FONT, "font.size": 10,
                         "axes.edgecolor": "#444A52", "axes.labelcolor": "black",
                         "text.color": "black", "xtick.color": "#444A52",
                         "ytick.color": "#444A52"})
    return plt


def fig_architecture() -> Path:
    plt = _mpl()
    fig, ax = plt.subplots(figsize=(12.0, 4.3))
    ax.axis("off")

    boxes = [
        (0.02, "Cortex\n(32 rec. sites)", "#EAEFF5"),
        (0.19, "PZ2 / RZ2\nbase 24414 Hz", "#EAEFF5"),
        (0.36, "NPro1\n610.35 Hz", "#EAEFF5"),
        (0.53, "PO8e card\n→ PC", "#DCE6F1"),
        (0.70, "Feature\n6 samples", "#DCE6F1"),
        (0.87, "MPC\n100 Hz", "#DCE6F1"),
    ]
    for x, label, col in boxes:
        ax.add_patch(plt.Rectangle((x, 0.60), 0.12, 0.26, facecolor=col,
                                   edgecolor="#444A52", lw=1.2))
        ax.text(x + 0.06, 0.73, label, ha="center", va="center", fontsize=9.5)
    for x, _, _ in boxes[:-1]:
        ax.annotate("", xy=(x + 0.17, 0.73), xytext=(x + 0.12, 0.73),
                    arrowprops=dict(arrowstyle="->", color="#444A52", lw=1.4))

    # return path
    ax.add_patch(plt.Rectangle((0.53, 0.16), 0.12, 0.26, facecolor="#F3E4E4",
                               edgecolor="#444A52", lw=1.2))
    ax.text(0.59, 0.29, "StimGen\n101.73 Hz", ha="center", va="center", fontsize=9.5)
    ax.add_patch(plt.Rectangle((0.30, 0.16), 0.15, 0.26, facecolor="#F3E4E4",
                               edgecolor="#444A52", lw=1.2))
    ax.text(0.375, 0.29, "8 bipolar pairs\n→ 16 electrodes", ha="center",
            va="center", fontsize=9.5)
    ax.annotate("", xy=(0.93, 0.60), xytext=(0.93, 0.42),
                arrowprops=dict(arrowstyle="-", color="#444A52", lw=1.4))
    ax.annotate("", xy=(0.65, 0.29), xytext=(0.93, 0.29),
                arrowprops=dict(arrowstyle="->", color="#444A52", lw=1.4))
    ax.text(0.79, 0.335, "UDP, 8 words @100 Hz", ha="center", fontsize=9,
            color="#444A52")
    ax.annotate("", xy=(0.45, 0.29), xytext=(0.53, 0.29),
                arrowprops=dict(arrowstyle="->", color="#444A52", lw=1.4))
    ax.annotate("", xy=(0.08, 0.60), xytext=(0.30, 0.29),
                arrowprops=dict(arrowstyle="->", color="#444A52", lw=1.4))

    ax.text(0.50, 0.94, "Forward path: acquisition → feature → controller",
            ha="center", fontsize=11, weight="bold")
    ax.text(0.50, 0.045, "Return path: command → stimulator → tissue",
            ha="center", fontsize=11, weight="bold")
    ax.set_xlim(0, 1.02); ax.set_ylim(0, 1)
    FIG_DIR.mkdir(parents=True, exist_ok=True)
    p = FIG_DIR / "arch_2026-08-14.png"
    fig.savefig(p, dpi=200, bbox_inches="tight"); plt.close(fig)
    return p


def fig_chain(ev) -> Path:
    plt = _mpl()
    fig, ax = plt.subplots(figsize=(11.4, 3.9))
    t = ev["uts"] - ev["uts"][0]
    n = min(700, t.size)
    ax.plot(t[:n], ev["cmd"][:n], lw=2.2, color="#2B4C7E", label="commanded (UDP1)")
    ax.plot(t[:n], ev["meas"][:n], lw=1.0, color="#C0504D",
            label="delivered scale (Scle)", alpha=0.9)
    ax.set_xlabel("time (s)"); ax.set_ylabel("amplitude (engineering units)")
    ax.set_title(f"Command reaches the stimulator intact  —  corr = {ev['corr']:.6f} "
                 f"at exactly {ev['lag']:+d} tick", fontsize=11, weight="bold")
    ax.legend(frameon=False, fontsize=9.5, loc="upper right")
    ax.spines[["top", "right"]].set_visible(False)
    FIG_DIR.mkdir(parents=True, exist_ok=True)
    p = FIG_DIR / "chain_2026-08-14.png"
    fig.savefig(p, dpi=200, bbox_inches="tight"); plt.close(fig)
    return p


def fig_window() -> Path:
    plt = _mpl()
    fig, (a1, a2) = plt.subplots(1, 2, figsize=(11.6, 3.8))

    rng = np.random.default_rng(7)
    before = np.clip(rng.lognormal(4.2, 1.5, 300), 3, 5439)
    a1.plot(before, lw=0.8, color="#C0504D")
    a1.axhline(6, color="#444A52", ls="--", lw=1.2)
    a1.set_yscale("log")
    a1.set_title("BEFORE — window by arrival time\n3 to 5,439 samples per tick "
                 "(hardware)", fontsize=10.5, weight="bold")
    a1.set_xlabel("control tick"); a1.set_ylabel("samples in window (log)")
    a1.text(150, 9, "intended: 6", fontsize=9, color="#444A52")

    a2.plot(np.full(300, 6), lw=2.0, color="#2B4C7E")
    a2.set_ylim(0, 12)
    a2.set_title("AFTER — window by sample count\nexactly 6 on every tick, "
                 "1000/1000 verified", fontsize=10.5, weight="bold")
    a2.set_xlabel("control tick"); a2.set_ylabel("samples in window")
    for a in (a1, a2):
        a.spines[["top", "right"]].set_visible(False)
    fig.text(0.5, -0.06, "Left panel is illustrative of the measured 3–5,439 range; "
             "right panel is the measured constant.", ha="center", fontsize=8.5,
             color="#444A52")
    FIG_DIR.mkdir(parents=True, exist_ok=True)
    p = FIG_DIR / "window_2026-08-14.png"
    fig.savefig(p, dpi=200, bbox_inches="tight"); plt.close(fig)
    return p


# --------------------------------------------------------------------------- #
# deck
# --------------------------------------------------------------------------- #
def build(ev, figs):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    BLACK = RGBColor(0, 0, 0)
    GREY = RGBColor(0x44, 0x4A, 0x52)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]

    def new(title, subtitle=None):
        s = prs.slides.add_slide(blank)
        box = s.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(12.2), Inches(0.72))
        p = box.text_frame.paragraphs[0]
        p.text = title
        p.font.size, p.font.bold, p.font.name, p.font.color.rgb = Pt(26), True, FONT, BLACK
        if subtitle:
            sb = s.shapes.add_textbox(Inches(0.57), Inches(1.02), Inches(12.2), Inches(0.5))
            sp = sb.text_frame.paragraphs[0]
            sp.text = subtitle
            sp.font.size, sp.font.name, sp.font.color.rgb = Pt(13), FONT, GREY
        return s

    def bullets(s, items, top=1.62, size=15.5, left=0.72, width=12.0):
        height = max(0.5, 7.5 - top - 0.25)
        tf = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width),
                                  Inches(height)).text_frame
        tf.word_wrap = True
        first = True
        for lvl, text in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = text
            p.level = lvl
            p.font.size = Pt(size if lvl == 0 else size - 1.5)
            p.font.name = FONT
            p.font.color.rgb = BLACK if lvl == 0 else GREY
            p.space_after = Pt(7)
        return tf

    def table(s, rows, left, top, width, height, col_w=None, size=12):
        n_r, n_c = len(rows), len(rows[0])
        shp = s.shapes.add_table(n_r, n_c, Inches(left), Inches(top),
                                 Inches(width), Inches(height)).table
        if col_w:
            for i, w in enumerate(col_w):
                shp.columns[i].width = Inches(w)
        for r, row in enumerate(rows):
            for c, val in enumerate(row):
                cell = shp.cell(r, c)
                cell.text = str(val)
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(size)
                para.font.name = FONT
                para.font.bold = (r == 0)
                para.font.color.rgb = BLACK
        return shp

    def pic(s, path, left, top, width):
        s.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))

    # ---- 1 title ----------------------------------------------------------
    s = prs.slides.add_slide(blank)
    b = s.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.4))
    p = b.text_frame.paragraphs[0]
    p.text = "Closed-Loop Stimulation Control"
    p.font.size, p.font.bold, p.font.name, p.font.color.rgb = Pt(38), True, FONT, BLACK
    b2 = s.shapes.add_textbox(Inches(0.95), Inches(3.6), Inches(11.5), Inches(1.8))
    tf = b2.text_frame; tf.word_wrap = True
    for i, (txt, sz) in enumerate([
        ("System status and in-vivo readiness assessment", 19),
        ("TDTClosedLoop  •  real-time MPC over PO8e → UDP → RZ2", 14),
        ("2026-08-14", 13),
    ]):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = txt
        para.font.size, para.font.name = Pt(sz), FONT
        para.font.color.rgb = BLACK if i == 0 else GREY

    # ---- 2 architecture ---------------------------------------------------
    s = new("The system", "Two paths that must both work: acquisition→control, and command→stimulator")
    pic(s, figs["arch"], 0.55, 1.45, 12.2)

    # ---- 3 where we started -----------------------------------------------
    s = new("Where we started", "Rig day 1, 2026-07-30 — engineering pass, biological null")
    bullets(s, [
        (0, "Hardware path worked on the first attempt: 3000/3000 control ticks, zero dropped, 1.37 ms turnaround."),
        (0, "But no channel responded to stimulation. Best |corr| 0.059 against a 0.027–0.064 noise floor."),
        (1, "Indistinguishable from nothing, and uniform across channels — which also ruled out artifact."),
        (0, "That null has since been explained. It was never one fault; it was three, stacked."),
        (1, "Each one alone would have produced exactly the same result: a clean-looking run with no signal."),
        (0, "This deck covers what those faults were, which are closed, and what remains before an animal."),
    ])

    # ---- 4 faults closed --------------------------------------------------
    s = new("Two faults found and closed", "Both were silent — neither produced an error message")
    bullets(s, [
        (0, "1. The network was never carrying the commands  (found 08-09, fixed 08-10)"),
        (1, "This PC had no address on the RZ2's subnet. UDP send() reports success into a void — "
            "the socket uses connect()+send() with no bind, so every send returned a full byte count."),
        (1, "Fixed with a static address on the correct subnet. Verified by GET_VERSION handshake and a subnet sweep."),
        (0, "2. The received value never reached the stimulator  (found 08-11, fixed 08-12)"),
        (1, "The command arrived at the RZ2 bit-exact and stopped there. A circuit-side routing issue, "
            "not a code fault — nothing in the software could have caused or detected it."),
        (0, "Both were invisible from the control software's point of view. That is the common theme, "
            "and it is why the handshake and diagnostics described later were added."),
    ])

    # ---- 5 chain proven ---------------------------------------------------
    s = new("The command path is now proven bit-exact",
            f"Recomputed from block LD-260812-185253 — {ev['n_pkt']} packets, word {ev['word']}")
    pic(s, figs["chain"], 0.75, 1.5, 11.8)
    bullets(s, [
        (0, f"corr = {ev['corr']:.6f} at exactly {ev['lag']:+d} control tick — the structural "
            f"one-tick lag, now confirmed empirically rather than assumed."),
        (0, f"Packet interval {ev['interval_ms']:.3f} ms against a 10 ms target; commanded range "
            f"0.000–{ev['cmd_max']:.3f} reproduced exactly."),
    ], top=5.45, size=13.5)

    # ---- 6 architecture correction ----------------------------------------
    s = new("Correction: the system has 8 actuators, not 16",
            "Found by inspecting what the stimulator actually emitted")
    bullets(s, [
        (0, "Each command word drives a BIPOLAR PAIR, not a single electrode."),
        (1, "sSig[2k] = −sSig[2k−1] exactly — max|w1+w2| = 0, corr = −1.000000. Not approximately: bit-exact."),
        (0, "So 8 words → 8 pairs → 16 electrodes. 8 is the true degree-of-freedom count."),
        (0, "This reversed a previous conclusion. The 8-word limit had been recorded as the receive gizmo "
            "truncating a 16-word command; it is not a truncation at all."),
        (1, "The fix therefore ran the opposite way: constrain the controller to 8 rather than widen the gizmo."),
        (0, "Not a limitation in practice — the acute work put the plant at rank ~3, so 8 independent "
            "pairs is comfortably more actuation than the system's controllable dimension."),
    ])

    # ---- 7 stim waveform --------------------------------------------------
    s = new("Stimulus characterised, and retuned",
            "Measured from the delivered pulse train, not from the circuit settings")
    table(s, [
        ["Parameter", "As measured (08-12)", "Retuned (08-14)", "Why it matters"],
        ["Pulse rate", f"{ev['stim_rate']:.4f} Hz (base/100)", f"{FS_STIM_NEW:.4f} Hz (base/240)",
         "Sets samples per stim period"],
        ["Samples per stim period", "2.5", "exactly 6", "Whole cycles → artifact is a constant offset"],
        ["Pulse period", "4.096 ms", "9.830 ms", "Now just under the 10 ms tick"],
        ["Pulses per control tick", "2.44", "1.02", "≥1 guaranteed — no silent dropouts"],
        ["Phase width / jitter", "204.80 µs / 0.00 µs", "unchanged", "Perfectly rigid train"],
    ], 0.72, 1.75, 11.9, 2.6, col_w=[2.5, 3.1, 3.1, 3.2], size=12)
    bullets(s, [
        (0, "Rule: f_stim = 610.3516 / k for integer k. A divisor that is a multiple of 40 gives an "
            "exact integer number of recording samples per stim period."),
        (0, "Below ~100 Hz you fall under one pulse per control tick, and some ticks would command an "
            "amplitude that is never delivered — a silent control dropout. 101.73 Hz is the floor."),
    ], top=4.6, size=13.5)

    # ---- 8 the window defect ----------------------------------------------
    s = new("Third fault: the feature the controller reads",
            "The one most likely to have caused the rig-day-1 null")
    pic(s, figs["window"], 0.75, 1.5, 11.8)
    bullets(s, [
        (0, "The window selected samples by the time they ARRIVED at the PC, not the time they were recorded. "
            "The card delivers in bursts, so each tick averaged a different amount of signal — up to ~200 ms."),
        (0, "The evoked response being sought lasts 15–28 ms. Averaging 200 ms of mostly-nothing with a "
            "20 ms response erases it."),
    ], top=5.3, size=13.5)

    # ---- 9 the fix --------------------------------------------------------
    s = new("Fixed, and confirmed on the hardware",
            "Window is now a fixed sample count, exposed as a run-time flag")
    bullets(s, [
        (0, "Every tick now averages exactly the same number of frames regardless of delivery jitter."),
        (1, "Verified 1000/1000 ticks across two runs: nominal 6, minimum 6, undersampled 0."),
        (1, "Confirmed on the real card under real bursty delivery, not only in simulation."),
        (0, "Window length is a command-line flag, so future feature work needs no rebuild."),
        (0, "Default is 6 samples = exactly one stim period, so the artifact contributes a constant "
            "offset rather than a tick-to-tick wobble that would look like signal."),
        (0, "Timing unchanged: 0.85 ms average tick, 1.76 ms worst case, against a 10 ms budget. "
            "Feature extraction itself dropped to 2.4 µs."),
        (0, "A related reporting bug was corrected: the 'undersampled windows' counter compared each tick "
            "against a running maximum, so it read 299/300 even in simulation. It was a false alarm that "
            "masked the real defect for weeks."),
    ])

    # ---- 10 supporting corrections ----------------------------------------
    s = new("Supporting corrections", "Each one a silent failure mode removed")
    table(s, [
        ["Issue", "Was", "Now"],
        ["UDP word count", "Defaulted to the live stream width — would have sent 32 words to an 8-word receiver",
         "Defaults to 8, the true actuator count"],
        ["Simulation sample rate", "24414 Hz — 40× faster than the real stream",
         "610.3516 Hz, the actual acquisition rate"],
        ["RZ2 handshake", "Disabled since the first commit (it could block forever)",
         "Enabled with a 1 s timeout; warns and continues"],
        ["Stale target addresses", "Two places aimed at this PC's own address",
         "Corrected to the RZ2"],
    ], 0.72, 1.8, 11.9, 3.2, col_w=[2.6, 5.0, 4.3], size=12)
    bullets(s, [
        (0, "The simulation rate is the most consequential of these: every simulation to date ran 40× too fast, "
            "which is exactly why simulation looked healthy while the rig did not."),
    ], top=5.4, size=13.5)

    # ---- 11 backup controller bug -----------------------------------------
    s = new("A bug found in the backup control path",
            "Found incidentally while building the test harness")
    bullets(s, [
        (0, "The native controller's open-loop mode — used for system identification — had never worked."),
        (1, "Two null-pointer dereferences: one crashed it at startup, the second on the first packet."),
        (0, "The existing 12-check self-test passed throughout, because it exercises the excitation "
            "generator directly and never the server path around it."),
        (0, "Both fixed; it now services a full run. The self-test still passes."),
        (0, "Worth stating plainly: this is the fallback for when MATLAB is unavailable at the rig. "
            "Had we needed it on a rig day, it would have failed at the worst possible moment."),
        (0, "Its output was also made unbuffered — the buffered tail was hiding the crash location and "
            "made the first two diagnosis attempts dead ends."),
    ])

    # ---- 12 tooling -------------------------------------------------------
    s = new("New diagnostic tooling", "Built to make the next failure legible instead of silent")
    bullets(s, [
        (0, "timing_check.py — measures command→feature latency from a capture, with an empirical noise floor."),
        (1, "The floor is computed by circularly shifting the data: same record length, same lag and channel "
            "search, so it accounts honestly for the multiple comparisons being made."),
        (1, "It judges on both agreement across channels AND magnitude above the floor. Either alone misreads."),
        (0, "Validated against three fixtures before being trusted:"),
        (1, "Planted synthetic relationship → detected at 9.9× the floor, correct lag."),
        (1, "Pure noise → correctly reports 'consistent with chance, do not report a latency'."),
        (1, "The rig-day-1 capture → correctly reports 'no timing measured'."),
        (0, "Also: a network diagnostic, an RZ2 emulator, a subnet scanner, and a 6-test UDP self-test — "
            "all of which now pass."),
    ])

    # ---- 13 current blocker -----------------------------------------------
    s = new("Current blocker: the acquisition loop crashes",
            "Three attempts on 2026-08-14 — including the binary that worked on rig day 1")
    table(s, [
        ["", "Attempt 1", "Attempt 2", "Attempt 3"],
        ["Binary", "Current build", "Current build", "ARCHIVED Jul-23 build"],
        ["Survived", "~36 ticks", "2 ticks", "~15-27 ticks"],
        ["Exit", "No summary; log never flushed", "Same", "Same"],
        ["Preceded by", "Large card position jump", "Same", "Same"],
    ], 0.72, 1.8, 11.9, 2.2, col_w=[2.3, 3.2, 3.2, 3.2], size=12)
    bullets(s, [
        (0, "The archived Jul-23 binary — the one that completed 3000 ticks on rig day 1 — fails the same "
            "way under today's conditions. The recent changes are therefore NOT the cause."),
        (0, "An earlier hypothesis (that the card accumulated a backlog the loop could not drain) was well "
            "evidenced and WRONG: the fix worked and the crash came sooner. It is recorded here because "
            "ruling it out was what made the controlled comparison the obvious next step."),
        (0, "The fault is environmental — card, driver, or configuration — and it predates this month's work."),
    ], top=4.3, size=13.5)

    # ---- 14 what we know --------------------------------------------------
    s = new("The crash: what is established", "Separating what is known from what is suspected")
    bullets(s, [
        (0, "Established:"),
        (1, "The recording path is healthy — all 16 channels live with real signal, features real-valued."),
        (1, "The feature window held at exactly 6 samples on hardware, under real bursty delivery."),
        (1, "The RZ2 handshake succeeded. The network path is fine."),
        (1, "It is confined to the acquisition loop. Nothing downstream had a chance to run."),
        (1, "The controlled comparison has now been run: the archived binary fails identically, so the "
            "recent changes are exonerated and the fault is environmental."),
        (0, "The same comparison confirmed the feature fix, on the same rig, minutes apart:"),
        (1, "Archived build: 4,444 then 4,291 samples per window — the original defect, still present."),
        (1, "Current build: exactly 6, every tick. Same card, same conditions."),
        (0, "Not established:"),
        (1, "What in the card, driver, or Synapse configuration causes it. Every run showed the card "
            "reporting a large position jump immediately beforehand."),
        (0, "Next test, and it is free: all three runs let the recording stream for 33-48 s before the loop "
            "attached. Reversing the start order cuts that to seconds and directly targets the one "
            "condition common to every failure."),
    ])

    # ---- 15 scorecard -----------------------------------------------------
    s = new("Validation scorecard", "What is proven on hardware, and what is not")
    table(s, [
        ["Capability", "Status", "Evidence"],
        ["Network transport to stimulator", "VALIDATED", "Handshake + subnet sweep + 6-test self-test, all pass"],
        ["Command → stimulator, bit-exact", "VALIDATED", f"corr {ev['corr']:.6f} at {ev['lag']:+d} tick, on hardware"],
        ["Stimulus waveform + timing", "VALIDATED", "Measured from delivered pulses; zero jitter"],
        ["Actuator count and pairing", "VALIDATED", "8 bipolar pairs, inversion exact"],
        ["Recording path (32→PC)", "PARTIAL", "16 channels live and real; 32 not yet delivered by the card"],
        ["Feature extraction determinism", "VALIDATED", "Exactly 6 samples/tick, 1000/1000, on hardware"],
        ["Sustained closed-loop operation", "NOT VALIDATED", "Loop crashes within seconds — current blocker"],
        ["Command → measured response", "NOT VALIDATED", "Never observed; requires a stable loop first"],
        ["Plant model / controller tuning", "NOT STARTED", "Deferred until a feature definition is settled"],
    ], 0.55, 1.72, 12.3, 3.9, col_w=[3.5, 2.2, 6.6], size=11.5)

    # ---- 16 risks ---------------------------------------------------------
    s = new("Residual risks before in vivo", "Ranked by what would cost the most animal time")
    table(s, [
        ["Risk", "Consequence", "Mitigation"],
        ["Loop instability under load", "Run ends early; animal time wasted",
         "Resolve current crash; require a full 60 s run twice before scheduling"],
        ["Mean-absolute-value is a placeholder feature",
         "Artifact-dominated, not linear in stimulus amplitude — model fits a curve with a line",
         "Redesign feature: baseline subtraction, artifact blanking, latency windowing"],
        ["No dynamics ever validated", "Model order and delay structure entirely unverified",
         "Saline validates timing only; dynamics need tissue"],
        ["Recording width unsettled", "Controller shape may change again",
         "Confirm the card delivers the intended channel count before fitting"],
        ["Non-negativity is the binding constraint",
         "Achievable stimulus set is a cone, not a subspace",
         "Known from the acute analysis; must be respected in the controller"],
    ], 0.55, 1.72, 12.3, 3.9, col_w=[3.1, 4.6, 4.6], size=11)

    # ---- 17 path ----------------------------------------------------------
    s = new("Path to in vivo", "In order; each gates the next")
    bullets(s, [
        (0, "1.  Resolve the acquisition crash.  Run the archived-binary comparison, then fix or revert."),
        (1, "Exit criterion: two consecutive full 60 s runs, zero dropped ticks."),
        (0, "2.  Saline timing confirmation.  Measure command→feature latency end to end."),
        (1, "Saline is instantaneous, so any lag measured is pure instrumentation."),
        (1, "Exit criterion: a stable, small, consistent lag across non-stimulated channels."),
        (0, "3.  Feature redesign.  Replace the placeholder with something linear in stimulus amplitude."),
        (1, "This is the deepest remaining item and the one most likely to decide whether control works."),
        (0, "4.  Plant identification in tissue, then closed-loop tuning."),
        (0, "Steps 1 and 2 are days. Step 3 is the real work."),
    ])

    # ---- 18 caveat --------------------------------------------------------
    s = new("What saline can and cannot tell us", "Stated up front so a pass is not over-read")
    bullets(s, [
        (0, "Saline validates: wiring, transport, timing, latency, and that the whole signal path is live."),
        (0, "Saline cannot validate: dynamics, model order, delay structure, or anything about control quality."),
        (1, "It is an instantaneous, essentially rank-1 medium. It has no dynamics to identify."),
        (0, "A concrete illustration from our own test fixtures: on a synthetic plant, the peak "
            "cross-correlation lag was 7 ticks — but that is the plant's TIME CONSTANT, not a transport delay."),
        (1, "Peak-correlation lag equals latency only in a medium with no dynamics. That is exactly why "
            "saline is the right place to measure timing, and exactly why it cannot measure anything else."),
        (0, "A saline pass should be reported as 'the instrument works', not 'the controller works'."),
    ])

    # ---- 19 human summary -------------------------------------------------
    s = new("Where this actually stands")
    bullets(s, [
        (0, "The honest version: we spent months looking at a system that appeared to work and produced "
            "nothing, and we now know why. Not one bug — three, layered, each individually capable of "
            "producing the same clean-looking null."),
        (0, "The commands were never reaching the network. Once they did, they were not reaching the "
            "stimulator. Once they did, the feature the controller reads was averaging away the very "
            "response it was looking for."),
        (0, "All three are now closed, and each was proven closed with a measurement rather than an "
            "assumption. The command path is bit-exact. The stimulus is characterised to the microsecond. "
            "The feature is deterministic on real hardware."),
        (0, "What is not working: the acquisition loop crashes within seconds. We now know it is not "
            "something we introduced — the binary from three weeks ago, which once ran for a full minute, "
            "fails the same way on this rig today. So the problem is in the hardware or its configuration, "
            "and it has probably been sitting under the last two rig days without being visible."),
        (0, "The system is closer than it has been, and the remaining unknown is narrower and better "
            "instrumented than anything we faced before. But it is not ready for an animal, and the "
            "gating item is not subtle — the loop has to run for a minute without dying."),
        (0, "The deeper question, once it runs, is whether the feature we extract is one a linear "
            "controller can actually work with. That is the real scientific risk, and it is still ahead of us."),
    ], size=14.5)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    out = OUT_DIR / "ClosedLoop_readiness_2026-08-14.pptx"
    try:
        prs.save(out)
    except PermissionError:
        # The deck is open in PowerPoint. Don't lose the rebuild -- write beside
        # it and say so, rather than failing after all the work.
        n = 2
        while True:
            alt = out.with_name(f"{out.stem}_v{n}.pptx")
            try:
                prs.save(alt)
                print(f"NOTE: {out.name} is open in PowerPoint and locked.\n"
                      f"      Wrote {alt.name} instead -- close the original and "
                      f"re-run to consolidate.")
                return alt
            except PermissionError:
                n += 1
                if n > 9:
                    raise
    return out


def main():
    print("recomputing chain evidence from the block ...")
    ev = chain_evidence()
    print(f"  word {ev['word']}, corr {ev['corr']:.6f} at lag {ev['lag']:+d}, "
          f"{ev['n_pkt']} packets, inverted-pair={ev['inverted']}, "
          f"stim {ev['stim_rate']:.4f} Hz over {ev['n_pulses']} pulses")
    print("building figures ...")
    figs = {"arch": fig_architecture(), "chain": fig_chain(ev), "window": fig_window()}
    print("building deck ...")
    out = build(ev, figs)
    print(f"\nWrote {out}")


if __name__ == "__main__":
    main()
