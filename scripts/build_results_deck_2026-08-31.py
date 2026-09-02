"""Results deck for the 2026-08-31 acute closed-loop experiment.

Formatting per standing preference: Arial throughout, BLACK headings (never
blue), 11pt body. One figure per content slide wherever one exists (professor
preference). Figures + numbers come from day_2026-08-31/analysis/ (the three
overnight analysis families); headline numbers quoted from their summary
reports; each family's findings list is read from its *_summary.json at build
time so the deck cannot drift from the analysis.

Output: PythonIntanAnalysis/outputs/Synthesis/AcuteClosedLoop_2026-08-31_results.pptx
Run with the PythonIntanAnalysis venv python.
"""
from __future__ import annotations

import json
from pathlib import Path

REPO = Path(__file__).resolve().parents[1]
ANA = REPO / "day_2026-08-31" / "analysis"
OUT_DIR = REPO.parent / "PythonIntanAnalysis" / "outputs" / "Synthesis"
FONT = "Arial"
BODY_PT = 11.0


def findings(name):
    try:
        d = json.loads((ANA / name).read_text())
        f = d.get("findings", [])
        return [str(x) for x in f]
    except Exception as e:  # deck still builds if a JSON is malformed
        return [f"({name} unreadable: {e})"]


def build():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    BLACK = RGBColor(0, 0, 0)
    GREY = RGBColor(0x44, 0x4A, 0x52)
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]

    def new(title, subtitle=None):
        s = prs.slides.add_slide(blank)
        tb = s.shapes.add_textbox(Inches(0.55), Inches(0.3), Inches(12.2), Inches(0.7))
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = title
        r.font.name = FONT; r.font.size = Pt(23); r.font.bold = True
        r.font.color.rgb = BLACK
        if subtitle:
            tb2 = s.shapes.add_textbox(Inches(0.57), Inches(0.95), Inches(12.2), Inches(0.45))
            r2 = tb2.text_frame.paragraphs[0].add_run()
            r2.text = subtitle
            r2.font.name = FONT; r2.font.size = Pt(12); r2.font.color.rgb = GREY
        return s

    def bullets(s, items, top=1.5, size=BODY_PT, left=0.72, width=12.0):
        tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width),
                                  Inches(max(0.5, 7.5 - top - 0.25)))
        tf = tb.text_frame; tf.word_wrap = True
        first = True
        for level, text in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.level = level
            r = p.add_run(); r.text = text
            r.font.name = FONT; r.font.size = Pt(size)
            r.font.color.rgb = BLACK if level == 0 else GREY
            p.space_after = Pt(5)
        return tb

    def fig(s, name, left, top, width):
        p = ANA / name
        if p.is_file():
            s.shapes.add_picture(str(p), Inches(left), Inches(top), Inches(width))
        else:
            bullets(s, [(0, f"[figure missing: {name}]")], top=top)

    def fig_slide(title, subtitle, name, notes=None, fw=10.6, fl=1.35, ft=1.35):
        s = new(title, subtitle)
        fig(s, name, fl, ft, fw)
        if notes:
            bullets(s, notes, top=6.55, size=10.0)
        return s

    # 1 -- title
    s = prs.slides.add_slide(blank)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.6))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = "Closed-Loop Biomimetic Microstimulation: First In-Vivo Tracking"
    r.font.name = FONT; r.font.size = Pt(34); r.font.bold = True; r.font.color.rgb = BLACK
    tb2 = s.shapes.add_textbox(Inches(0.82), Inches(3.5), Inches(11.7), Inches(1.2))
    r2 = tb2.text_frame.paragraphs[0].add_run()
    r2.text = ("Acute experiment 2026-08-31 -- 32-ch cortical recording, thalamic stim "
               "(2 bipolar pairs), MPC closed loop vs Choi-2016 open loop on paired "
               "randomized interleaved touch targets. 21 blocks; every run 22,200/22,200 "
               "ticks, zero drops.")
    r2.font.name = FONT; r2.font.size = Pt(14); r2.font.color.rgb = GREY

    # 2 -- day at a glance
    s = new("The day in five acts")
    bullets(s, [
        (0, "1. 32-channel migration proven live: banner exact(32/32) on every run; touch "
            "battery 10/10 blocks, split-half 0.92-0.99, peaks 382-709 uV, modal best ch 8."),
        (0, "2. A NO-GO reversed: all standard fitters refused (sweep best |corr| 0.071). "
            "Raw-LFP signed averages found the response the fitters missed -- recruitment "
            "threshold ~13-18 uA; sub-threshold amps had drowned the linear fits."),
        (0, "3. Operating-point identification: PRBS 10-30 uA held above the knee on pairs "
            "1+4 -> |corr| 0.22-0.25 on ch 8 -> model (2-in/1-out, order 4, uOff~20 uA)."),
        (0, "4. Three volt-scale numerics bugs found and fixed live (Choi mu, Choi horizon "
            "tiling, MPC QP below solver tolerance -> QWeight 1e12 / RWeight 100)."),
        (0, "5. All four arm runs TRACKING: MPC rBest 0.712-0.725 at a stable +2-tick "
            "(~20 ms) lag; Choi 0.704-0.716 with drifting lag; drift re-probe shows the "
            "plant itself held."),
    ])

    # 3 -- touch battery
    fig_slide("The targets: 10-site touch battery at 32 channels",
              "150 contacts/site via the nThw thwacker; templates are the references the "
              "controllers must reproduce",
              "spat_templates.png",
              [(0, "9/9 real sites 380-710 uV, peak 25-29 ms, best ch 8 (LP: ch 6); SHAM flat "
                   "(40 uV). Touch site identity is decodable from cortex: 48% ten-way "
                   "Mahalanobis classification vs 10% chance.")], fw=9.8, fl=1.75)

    # 4 -- the pivot
    fig_slide("The pivot: a repositioning averted by finding the recruitment threshold",
              "Signed raw-LFP averages on data the fitters had refused",
              "spat_recruitment.png",
              [(0, "Pairs 4 and 1 drive S1 (+241 / +121 uV @ 13 ms, z = 23.5 / 9.7) with the "
                   "touch spatial footprint (r = 0.96-0.99). The knee at ~13-18 uA is why "
                   "zero-origin fits failed: lesson -- signed-average the raw LFP before "
                   "moving an electrode.")], fw=9.8, fl=1.75)

    # 5 -- event gallery
    fig_slide("Event-triggered tracking: all five sites, all arms",
              "Achieved ch-8 feature vs target, per site x arm-run; Hold = the accidental "
              "tonic-only control",
              "sci_event_gallery.png",
              [(0, "MPC event-average vs target r = 0.88-0.95 per site; Choi 0.66-0.90; the "
                   "Hold control is flat (r 0.004 +/- 0.125) -- tracking is controller "
                   "action, not artifact of the analysis.")], fw=10.2, fl=1.55)

    # 6 -- per-site scores
    fig_slide("Per-event scores by site", "Per-event r (lag-corrected) and peak ratio",
              "sci_persite_scores.png",
              [(0, "Caveat (honest): both arms overshoot single-tick peaks ~2x (median peak "
                   "ratio 2.26 MPC / 1.97 Choi) with onset jitter; event-averaged peaks "
                   "match. Targets below the spontaneous floor (SHAM's 0.038 mV bump) are "
                   "unreachable by either arm.")])

    # 7 -- paired comparison
    fig_slide("MPC vs Choi, paired on identical schedules: a tie on shape -- with a twist",
              "Same seeded event sequence in both arms; per-event paired differences",
              "sci_paired.png",
              [(0, "Raw scoring at each arm's own best lag reads as a TIE -- but with "
                   "artifact ticks (0-2 post-onset) EXCLUDED (ctl_ redo, 09-01) MPC wins: "
                   "pooled dr +0.079 [0.046, 0.112], p<1e-4, 65% of events -- driven by r1 "
                   "(+0.163; Choi's r1 parity was partly artifact samples), r2 a true tie. "
                   "Plus MPC's latency STABILITY (~2.1 ticks both runs vs Choi drifting "
                   "1.3->1.8) and within-run stability (next slide). Both arms sit far "
                   "above the Hold noise floor (r 0.06 vs ~0.5-0.6) -- tracking is real.")])

    # 8 -- time course
    fig_slide("The value of feedback appears over time",
              "Per-event tracking vs event number; the plant itself held (re-probe |corr| "
              "0.224 -> 0.248)",
              "sci_timecourse.png",
              [(0, "Choi decays within-run even at per-event best lag (0.81 -> 0.70, slope "
                   "-0.13 r per 100 events, significant, both runs); MPC never degrades and "
                   "improves in run 2. Extrapolation: MPC wins outright on runs longer than "
                   "~100 events. (Drift-check caveat: the final re-probe ran with degraded "
                   "PC timing -- 8.5% dropped ticks -- treat its fit as indicative only.)")])

    # 9 -- SHAM
    fig_slide("SHAM catch trials: the controllers do not invent touches",
              "Interleaved sham events, same scoring as real sites",
              "sci_sham.png",
              [(0, "False-touch rate MPC 0/32, Choi 1/33 (Hold base rate 1/16); real-event "
                   "hit rate 95% / 86%; d-prime 4.2 / 3.6. Controller sham modulation sits "
                   "at or below the spontaneous floor.")])

    # 10 -- spatial footprints
    fig_slide("Spatial footprints: probes match touch; arm averages carry artifact",
              "32-channel profiles -- touch templates vs single-pulse probes vs arm events",
              "spat_footprints.png",
              [(0, "Single pulses on pairs 1/4 reproduce the touch footprint at r = "
                   "0.96-0.99. BUT raw-LFP averages during the arms are contaminated by "
                   "overlapping continuous-stim artifact (Choi worst: 560-810 uV profiles "
                   "anticorrelated with touch). Online tracking lives in the processed "
                   "feature; biological-space confirmation during tonic stim needs "
                   "artifact-aware analysis (queued: pulse-resolved 24 kHz, signed "
                   "features, trim).")])

    # 11 -- Mahalanobis
    fig_slide("Site identity: the honest negative result",
              "Mahalanobis distance of arm-evoked patterns to each touch site's "
              "trial distribution",
              "spat_mahalanobis.png",
              [(0, "Stim-for-site does NOT land nearest its own touch site: 5-way accuracy "
                   "19% (MPC) / 22% (Choi) vs 20% chance. Expected from theory: one control "
                   "channel + a fixed stim footprint cannot SELECT which site it reproduces "
                   "(the acute rank/selectivity ceiling). The fix is the standing goal: "
                   "multichannel control at maximum achievable rank, with more pairs.")])

    # 12 -- charge + strategy
    fig_slide("Same fidelity, 30% less charge: the controllers' strategies differ",
              "Command traces and charge ledger on identical schedules",
              "eng_u_strategy.png",
              [(0, "Choi rides the 30 uA cap 48% of ticks and slews 4.9x harder; MPC stays "
                   "within +/-0.5 uA of hold half the time and never fully shuts off. "
                   "Choi delivered +29.5% more charge per run for tied fidelity. Tonic hold "
                   "is 97% of MPC's charge total (see eng_charge.png in the analysis "
                   "folder).")])

    # 13 -- engineering
    fig_slide("Engineering: the loop earned the day",
              "Latency distributions and tick health across all runs",
              "eng_latency.png",
              [(0, "MATLAB MPC server p50 1.69 ms (p99 29 ms; 5-6.5% loop timeouts absorbed "
                   "by hold-last, freshTicks >= 99.2%); cpp server p50 17 us -- ~100x. All "
                   "four arm runs: p95 tick error < 1.9 ms, zero dropped control ticks, no "
                   "mid-run PLL resyncs.")])

    # 14 -- findings lists (from the three JSONs, verbatim)
    for name, label in (("sci_summary.json", "Science-core findings (verbatim from analysis)"),
                        ("spat_summary.json", "Spatial/advanced findings (verbatim)"),
                        ("eng_summary.json", "Engineering findings (verbatim)")):
        f = findings(name)
        s = new(label, name)
        bullets(s, [(0, t) for t in f[:12]], size=10.5)

    # 15 -- good news / bad news (updated 09-01 with the artifact-aware redo)
    s = new("Summary: good news / bad news (artifact-aware, 09-01)")
    bullets(s, [
        (0, "GOOD: closed-loop biomimetic tracking works and survives artifact cleaning "
            "-- both arms sit far above the Hold noise floor (per-event r ~0.5-0.6 vs "
            "0.06); stable ~20 ms latency; d-prime ~4 sham rejection; 30% charge savings."),
        (0, "GOOD (upgraded): with artifact ticks excluded MPC BEATS Choi (pooled dr "
            "+0.079, p<1e-4; r1 +0.163, r2 tie) on top of its stability advantages -- "
            "Choi's apparent parity was partly artifact samples."),
        (0, "GOOD: no order effects in the interleaved design (residual autocorr null "
            "in all 5 runs) -- the randomized schedule did its job; paired stats valid."),
        (0, "BAD (artifact-robust): no site-selectivity in space (cleaned 23%/22% vs 20% "
            "chance) OR time course (intent decoding null). One control channel + one "
            "footprint cannot choose WHICH touch it reproduces. MIMO + more pairs is the path."),
        (0, "BAD: per-event amplitude transfer is absent in the cleaned estimate (slope "
            "~0.1 vs 1; single-event SNR-limited) -- event-averaged tracking is real, "
            "per-event grading is not demonstrated. Choi's step transient (~2 mV "
            "polarization at command steps) needs a slew penalty next time."),
        (0, "NEXT: MIMO/max-rank control with all 8 pairs (9/3); Choi --lam smoothness; "
            "signed/trimmed features; longer runs to cash in MPC stability; 64-ch map."),
    ])

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    out = OUT_DIR / "AcuteClosedLoop_2026-08-31_results.pptx"
    try:
        prs.save(out)
    except PermissionError:
        out = out.with_name(out.stem + "_v2.pptx")
        prs.save(out)
    print(f"Wrote {out} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)"
          if False else f"Wrote {out}")


if __name__ == "__main__":
    build()
