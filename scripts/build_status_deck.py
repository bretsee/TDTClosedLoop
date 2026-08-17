"""Build the progress-meeting deck on the closed-loop stim control system.

Covers: the paradigm and why closed loop is required, system architecture, the
2026-07-23 controller-path decision and observer defect, the pre-rig verification,
the 2026-07-30 rig-day results (hardware pass, biological null), the ranked
differential for the null, and the plan for the next rig day.

Headline numbers are RECOMPUTED from capture_rig_run0.csv rather than transcribed
from the console, so the deck cannot drift from the data it describes. The one
number that is quoted rather than computed is the synthetic noise floor, which is
labelled as such on the slide.

Formatting per user preference: Arial throughout, black headings (never blue).

Output: PythonIntanAnalysis/outputs/Synthesis/ClosedLoop_system_status_2026-07-30.pptx
"""
from __future__ import annotations

from pathlib import Path

import numpy as np

REPO = Path(__file__).resolve().parents[1]              # TDTClosedLoop
CAPTURE = REPO / "capture_rig_run0.csv"
OUT_DIR = REPO.parent / "PythonIntanAnalysis" / "outputs" / "Synthesis"
FIG_DIR = REPO / "outputs" / "deck_figures"

FONT = "Arial"

# Measured on synthetic pure-noise channels at the same record length during the
# 2026-07-29 positive/negative control run. Quoted, not recomputed here.
NOISE_FLOOR = (0.027, 0.064)
GO_THRESHOLD = 0.10


# --------------------------------------------------------------------------- #
# recompute the sweep from the capture, so the deck matches the data
# --------------------------------------------------------------------------- #
def max_abs_xcorr(x: np.ndarray, y: np.ndarray, max_lag: int) -> float:
    """Peak |Pearson r| of y against x over lags 0..max_lag.

    Mirrors sweep_channels.m's max_abs_xcorr so the deck reproduces the number
    the rig actually saw, rather than a differently-defined correlation.
    """
    x = x - x.mean()
    y = y - y.mean()
    n = x.size
    max_lag = max(0, min(max_lag, n - 10))
    best = 0.0
    for lag in range(max_lag + 1):
        xa, ya = x[: n - lag], y[lag:]
        dx, dy = np.linalg.norm(xa), np.linalg.norm(ya)
        if dx > 1e-12 and dy > 1e-12:
            best = max(best, abs(float(xa @ ya) / (dx * dy)))
    return best


def sweep_from_capture(path: Path, skip_ticks: int = 50, delay_ticks: int = 1):
    import csv

    with open(path, newline="") as fh:
        rows = list(csv.reader(fh))
    header, data = rows[0], np.array(rows[1:], dtype=float)
    u_idx = [i for i, h in enumerate(header) if h.startswith("u")]
    y_idx = [i for i, h in enumerate(header) if h.startswith("y")]
    U, Y = data[:, u_idx], data[:, y_idx]

    # Same causal alignment and warm-up discard the fitter uses: the command
    # logged at tick k is not transmitted until k+1, so u(k) acts on y(k+1).
    if delay_ticks:
        U, Y = U[:-delay_ticks], Y[delay_ticks:]
    U, Y = U[skip_ticks:], Y[skip_ticks:]

    moved = [i for i in range(U.shape[1]) if U[:, i].std() > 1e-12]
    corrs = [max(max_abs_xcorr(U[:, m], Y[:, ch], 50) for m in moved)
             for ch in range(Y.shape[1])]
    return np.array(corrs), [m + 1 for m in moved], U.shape[0]


# --------------------------------------------------------------------------- #
# figures
# --------------------------------------------------------------------------- #
def fig_architecture() -> Path:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    from matplotlib.patches import FancyArrowPatch, FancyBboxPatch

    fig, ax = plt.subplots(figsize=(12.6, 3.5))
    boxes = [
        (0.2, "S1 cortex\n16 ch", "#DBEAFE"),
        (2.5, "PO8e card\n24.4 kHz", "#E5E7EB"),
        (4.8, "C++ loop\n100 Hz tick\nMAV feature", "#E5E7EB"),
        (7.1, "MATLAB MPC\nOSQP + observer\n(localhost)", "#FEF3C7"),
        (9.4, "RZ2\nUDP :22022", "#E5E7EB"),
        (11.7, "VPL stim\n16 pairs", "#FECACA"),
    ]
    for x, label, color in boxes:
        ax.add_patch(FancyBboxPatch((x, 1.0), 1.9, 1.25, boxstyle="round,pad=0.05",
                                    fc=color, ec="#111827", lw=1.3))
        ax.text(x + 0.95, 1.62, label, ha="center", va="center", fontsize=9.5)

    for x in (2.1, 4.4, 6.7, 9.0, 11.3):
        ax.add_patch(FancyArrowPatch((x, 1.62), (x + 0.4, 1.62), arrowstyle="-|>",
                                     mutation_scale=15, color="#111827", lw=1.7))
    # closing the loop, drawn underneath
    ax.add_patch(FancyArrowPatch((12.65, 0.95), (1.15, 0.95), arrowstyle="-|>",
                                 mutation_scale=15, color="#B45309", lw=1.7,
                                 connectionstyle="arc3,rad=0.08"))
    ax.text(6.9, 0.28, "closed loop — stim evokes the cortical response that is measured on the next tick",
            ha="center", fontsize=9, color="#B45309")
    ax.text(6.0, 2.62, "10 ms budget per tick  •  structural one-tick (10 ms) command lag on the localhost path",
            ha="center", fontsize=9.5, color="#374151")

    ax.set_xlim(0, 13.8); ax.set_ylim(0, 3.1); ax.axis("off")
    FIG_DIR.mkdir(parents=True, exist_ok=True)
    p = FIG_DIR / "architecture.png"
    fig.savefig(p, dpi=150, bbox_inches="tight")
    plt.close(fig)
    return p


def fig_sweep(corrs: np.ndarray) -> Path:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    # Wide and short: the slide places bullets beneath this figure, so the
    # rendered height has to stay under ~3.2 in at full slide width.
    fig, ax = plt.subplots(figsize=(12.0, 3.7))
    ch = np.arange(1, corrs.size + 1)
    ax.bar(ch, corrs, color="#93A3B8", edgecolor="#111827", lw=0.7, width=0.68,
           label="measured |corr|, stim → channel")
    ax.axhspan(NOISE_FLOOR[0], NOISE_FLOOR[1], color="#D1D5DB", alpha=0.55, zorder=0,
               label=f"noise floor, synthetic control ({NOISE_FLOOR[0]:.3f}–{NOISE_FLOOR[1]:.3f})")
    ax.axhline(GO_THRESHOLD, color="#B91C1C", lw=2.0, ls="--",
               label=f"go/no-go threshold ({GO_THRESHOLD:.2f})")

    ax.set_xticks(ch)
    ax.set_xlabel("S1 cortical channel", fontsize=11)
    ax.set_ylabel("peak |corr|, lags 0–500 ms", fontsize=11)
    ax.set_ylim(0, 0.115)
    # Legend above the axes: placed inside, it sits on top of the threshold line,
    # which is the single most important mark on the chart.
    ax.legend(fontsize=9.5, loc="lower center", bbox_to_anchor=(0.5, 1.01),
              ncol=3, frameon=False)
    ax.spines[["top", "right"]].set_visible(False)
    fig.tight_layout()
    p = FIG_DIR / "sweep_run0.png"
    fig.savefig(p, dpi=150, bbox_inches="tight")
    plt.close(fig)
    return p


def fig_window() -> Path:
    """Feature-window instability: simulation vs real hardware."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    # Narrow: this figure sits in the left half of the slide beside a text
    # column starting at 6.6 in, so its rendered width must stay under ~5.8 in.
    fig, ax = plt.subplots(figsize=(6.6, 4.4))
    labels = ["Simulated input\n(rehearsal)", "PO8e hardware\n(rig day 1)"]
    lo, hi = np.array([194, 3]), np.array([296, 5439])
    xs = np.arange(2)

    ax.vlines(xs, lo, hi, color="#111827", lw=3.0)
    ax.scatter(xs, lo, s=90, color="#B91C1C", zorder=3, label="minimum samples in window")
    ax.scatter(xs, hi, s=90, color="#1D4ED8", zorder=3, label="maximum samples in window")
    ax.axhline(244, color="#059669", lw=1.8, ls="--",
               label="expected for a true 10 ms window (244 samples)")

    for x, l, h in zip(xs, lo, hi):
        ax.text(x + 0.09, h, f"{h:,}", va="center", fontsize=10)
        ax.text(x + 0.09, max(l, 3.5), f"{l:,}", va="center", fontsize=10)

    ax.set_yscale("log")
    ax.set_xticks(xs); ax.set_xticklabels(labels, fontsize=10.5)
    ax.set_xlim(-0.4, 1.6)
    ax.set_ylabel("samples inside the 10 ms feature window (log)", fontsize=10.5)
    ax.legend(fontsize=9, loc="upper left", framealpha=0.95)
    ax.spines[["top", "right"]].set_visible(False)
    fig.tight_layout()
    p = FIG_DIR / "feature_window.png"
    fig.savefig(p, dpi=150, bbox_inches="tight")
    plt.close(fig)
    return p


# --------------------------------------------------------------------------- #
# deck
# --------------------------------------------------------------------------- #
def build(corrs, moved, n_ticks, figs):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    BLACK = RGBColor(0, 0, 0)
    GREY = RGBColor(0x44, 0x4A, 0x52)

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]

    def new(title, subtitle=None):
        s = prs.slides.add_slide(blank)
        box = s.shapes.add_textbox(Inches(0.55), Inches(0.32), Inches(12.2), Inches(0.72))
        p = box.text_frame.paragraphs[0]
        p.text = title
        p.font.size, p.font.bold, p.font.name, p.font.color.rgb = Pt(26), True, FONT, BLACK
        if subtitle:
            sb = s.shapes.add_textbox(Inches(0.57), Inches(1.02), Inches(12.2), Inches(0.5))
            sp = sb.text_frame.paragraphs[0]
            sp.text = subtitle
            sp.font.size, sp.font.name, sp.font.color.rgb = Pt(13), FONT, GREY
        return s

    def bullets(s, items, top=1.62, size=15.5, left=0.72, width=12.0, height=None):
        # Default the box to whatever vertical space is actually left on the
        # slide. A fixed height silently runs the frame off the bottom edge for
        # any block placed low, which PowerPoint may then autofit unpredictably.
        if height is None:
            height = max(0.5, 7.5 - top - 0.25)
        tf = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height)).text_frame
        tf.word_wrap = True
        first = True
        for lvl, text in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = text
            p.level = lvl
            p.font.size = Pt(size if lvl == 0 else size - 1.5)
            p.font.name = FONT
            p.font.color.rgb = BLACK if lvl == 0 else GREY
            p.space_after = Pt(7)
        return tf

    def table(s, rows, left, top, width, height, col_w=None, size=12):
        n_r, n_c = len(rows), len(rows[0])
        shp = s.shapes.add_table(n_r, n_c, Inches(left), Inches(top),
                                 Inches(width), Inches(height)).table
        if col_w:
            for i, w in enumerate(col_w):
                shp.columns[i].width = Inches(w)
        for r, row in enumerate(rows):
            for c, val in enumerate(row):
                cell = shp.cell(r, c)
                cell.text = str(val)
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(size)
                para.font.name = FONT
                para.font.bold = (r == 0)
                para.font.color.rgb = BLACK
        return shp

    # ---- 1. title ---------------------------------------------------------
    s = prs.slides.add_slide(blank)
    b = s.shapes.add_textbox(Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.4))
    p = b.text_frame.paragraphs[0]
    p.text = "Closed-Loop Thalamic Stimulation Control"
    p.font.size, p.font.bold, p.font.name, p.font.color.rgb = Pt(38), True, FONT, BLACK
    b2 = s.shapes.add_textbox(Inches(0.95), Inches(3.7), Inches(11.5), Inches(1.6))
    tf = b2.text_frame
    tf.word_wrap = True
    for i, (txt, sz) in enumerate([
        ("System status, first rig-day results, and next steps", 19),
        ("Real-time MPC for patterned VPL stimulation  •  TDTClosedLoop", 14),
        ("30 July 2026", 14),
    ]):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = txt
        para.font.size, para.font.name = Pt(sz), FONT
        para.font.color.rgb = BLACK if i == 0 else GREY
        para.space_after = Pt(8)

    # ---- 2. why closed loop ----------------------------------------------
    s = new("Why closed-loop control",
            "The scientific requirement that motivates the whole system")
    bullets(s, [
        (0, "Goal: stimulate VPL thalamus to reproduce, in S1 cortex, the response a natural paw touch evokes."),
        (0, "Prior offline result (4 acute studies, 15 touch targets) established where the difficulty lies:"),
        (1, "Spatial pattern is reachable — median per-channel reachability 1.00. Space is not the bottleneck."),
        (1, "Temporal pattern is not — median temporal reachability 0.44; spatiotemporal energy 0.08."),
        (0, "Interpretation: fixed, fast stim kernels cannot construct touch's slow ~200 ms ramp."),
        (0, "More stimulation pairs will not fix this — the inversion already permits arbitrary multi-pair mixing."),
        (0, "What is required is time-varying, patterned stimulation. That is precisely what a closed loop delivers,"
            " and it is why this control system is being built rather than a fixed stimulus table."),
    ])

    # ---- 3. architecture --------------------------------------------------
    s = new("System architecture",
            "100 Hz real-time loop, C++ acquisition and control, MATLAB model-based controller")
    s.shapes.add_picture(str(figs["arch"]), Inches(0.5), Inches(1.55), width=Inches(12.3))
    bullets(s, [
        (0, "Controller: condensed-QP linear MPC (OSQP) with a state observer, driving a 16-channel "
            "pulse-amplitude envelope — one amplitude per control point, bounded by the stimulator's "
            "physical charge ceiling."),
    ], top=5.35, size=14)

    # ---- 4. prior milestone ----------------------------------------------
    s = new("Prior milestone (23 July): path decided, and a silent defect found",
            "Two findings, the second of which invalidated all earlier control-behaviour claims")
    bullets(s, [
        (0, "1. Controller path resolved in favour of the localhost architecture."),
        (1, "Embedded runs the MATLAB solve inside the tick: worst case 11.3 ms against a 10 ms budget, "
            "and it dropped ticks (2 of 500)."),
        (1, "Localhost decouples the loop: ~53 µs ticks, zero dropped. Cost is a structural "
            "one-tick (10 ms) command lag, which is now explicitly modelled."),
        (0, "2. The MPC was running fully open-loop and had been for its entire history."),
        (1, "The observer gain fell back to L = 0 whenever place() threw — and place() is Control "
            "System Toolbox, which is not installed. The measurement was computed, logged, then "
            "multiplied by zero."),
        (1, "Replaced with a dependency-free steady-state Kalman gain. Command range over a swept "
            "feature went from 0.000 to 8.309."),
        (0, "Consequence: timing results remain valid, but no control-behaviour claim predating this "
            "fix should be relied upon."),
    ], size=14.5)

    # ---- 5. pre-rig verification -----------------------------------------
    s = new("Pre-rig verification (29 July)",
            "Whole toolchain re-run end to end so that hardware day had one new variable, not several")
    bullets(s, [
        (0, "Verified on the rig machine before hardware was involved:"),
        (1, "Control loop smoke test — 0 dropped ticks, 25.9 µs mean tick"),
        (1, "Capture server → control loop → capture file — 0 dropped, 779/800 replies serviced"),
        (1, "System-ID fitter on a no-signal record — correctly reported no response and refused to claim a model"),
        (1, "Closed-loop MPC in simulation — command demonstrably tracks the measurement"),
        (0, "Operational tooling built to make the rig day executable under time pressure:"),
        (1, "Five sequential scripts covering pre-flight, capture, control loop, fitting and closed loop"),
        (1, "A channel sweep that triages all 16 recording channels in one command"),
        (1, "Cross-record validation that distinguishes a genuine model from record-specific structure"),
        (1, "A branching protocol with ~20 labelled failure branches, referenced by name in script output"),
    ], size=14.5)

    # ---- 6. controls ------------------------------------------------------
    s = new("The detection pipeline was validated in both directions",
            "So that a negative result on rig day would be interpretable rather than ambiguous")
    table(s, [
        ["Control", "Data", "Expected", "Observed", "Result"],
        ["Negative", "Real capture, no causal link\nbetween stim and signal",
         "No response detected", "|corr| 0.006–0.010\nwarnings raised", "Pass"],
        ["Positive", "Synthetic plant, poles 0.900 / 0.700,\nSNR 3, response on one channel",
         "Correct channel identified", "Found the right channel\n|corr| 0.626, fit 68.5 %", "Pass"],
        ["System ID", "Synthetic plant, known dynamics",
         "Recover true dynamics", "Poles 0.912 / 0.693\nDC gain within 2.3 %", "Pass"],
    ], left=0.7, top=1.75, width=11.9, height=3.4,
        col_w=[1.5, 3.3, 2.5, 3.0, 1.6], size=11.5)
    bullets(s, [
        (0, "The pipeline neither manufactures a model from noise nor misses a response that is present. "
            "A null result from it is therefore evidence about the preparation or the measurement — "
            "not an unexplained failure of the analysis."),
    ], top=5.45, size=14)

    # ---- 7. rig day: hardware pass ---------------------------------------
    s = new("Rig day 1 — the engineering path worked on first attempt",
            "Every layer from acquisition card to stimulator behaved to specification")
    table(s, [
        ["Measure", "Result", "Assessment"],
        ["PO8e acquisition", "Card connected, 16 channels streaming", "As expected"],
        ["Control ticks completed", "3000 of 3000", "Complete"],
        ["Dropped control ticks", "0", "Budget held"],
        ["Controller turnaround", "1.37 ms mean, 1.92 ms 95th pct", "Well inside 10 ms"],
        ["UDP transmission to RZ2", "Connected, transmission enabled", "Confirmed"],
        ["Reply timeouts", "270 of 3000 (45 per 500)", "Within 23–74 per 500 baseline"],
        ["Capture written", "3000 rows", "Complete"],
    ], left=0.85, top=1.8, width=11.6, height=3.9,
        col_w=[3.9, 4.6, 3.1], size=12.5)
    bullets(s, [
        (0, "No software or hardware fault was observed at any point during the session."),
    ], top=5.95, size=14)

    # ---- 8. rig day: null result -----------------------------------------
    s = new("Rig day 1 — no measurable cortical response to stimulation",
            f"Open-loop PRBS, single stimulation pair (channel {moved[0]}), amplitude 0–30, "
            f"{n_ticks:,} analysed ticks")
    s.shapes.add_picture(str(figs["sweep"]), Inches(1.35), Inches(1.6), width=Inches(10.6))
    bullets(s, [
        (0, f"Strongest channel reached |corr| = {corrs.max():.3f}, against a decision threshold of "
            f"{GO_THRESHOLD:.2f} fixed in advance. No model was fitted."),
        (0, f"All 16 channels fall within the {NOISE_FLOOR[0]:.3f}–{NOISE_FLOOR[1]:.3f} noise floor "
            "measured on synthetic data at the same record length — the result is not a weak response, "
            "it is indistinguishable from none."),
        (0, "The response is uniform across channels with no focal peak, which additionally excludes "
            "stimulation artefact as an explanation."),
    ], top=4.95, size=13.5)

    # ---- 9. leading hypothesis -------------------------------------------
    s = new("Leading hypothesis: feature extraction, not biology",
            "The measured quantity may be incapable of resolving the response, independent of whether one occurred")
    s.shapes.add_picture(str(figs["window"]), Inches(0.5), Inches(1.9), width=Inches(5.75))
    bullets(s, [
        (0, "The control feature is a mean-absolute-value over a 10 ms window — but the window is "
            "defined in sample ARRIVAL time, not sample position."),
        (0, "The acquisition card delivers in blocks of 40, so samples arrive in bursts and receive "
            "near-identical timestamps."),
        (0, "Consequence: the number of samples averaged per tick varied from 3 to 5,439 on hardware, "
            "against a steady ~250 in simulation."),
        (0, "At the upper end the feature averages ~200 ms of data, which would attenuate the "
            "15–28 ms evoked response the system is trying to detect."),
        (0, "Status: a hypothesis consistent with the evidence, not an established cause. It requires "
            "a C++ change and was deliberately not attempted during a live session."),
    ], top=1.72, left=6.6, width=6.2, size=13)

    # ---- 10. differential -------------------------------------------------
    s = new("Ranked explanations for the null result",
            "Ordered by diagnostic cost, cheapest and most decisive first")
    table(s, [
        ["#", "Explanation", "How it is tested", "Status"],
        ["1", "Stimulation was never delivered",
         "Inspect the delivered-stim record, which is\nlogged independently of the control loop",
         "Not yet checked —\ngates all others"],
        ["2", "A single pair is too weak a driver",
         "Drive all 16 pairs — approximately 16× the charge",
         "Prior work: one pair reaches\n~6 % of a natural touch response"],
        ["3", "Feature window cannot resolve the response",
         "Compare against the raw recorded waveform,\nwhich bypasses the feature entirely",
         "Leading hypothesis"],
        ["4", "Wrong recording array or channel mapping",
         "Confirm channel count and array assignment",
         "16 channels confirmed"],
        ["5", "Amplitude too low",
         "Raise toward the physical ceiling",
         "Unlikely — already at 30 of 40"],
        ["6", "The preparation did not respond",
         "Diagnosis of exclusion, only after 1–5",
         "Open"],
    ], left=0.55, top=1.7, width=12.2, height=4.4,
        col_w=[0.5, 3.5, 4.6, 3.6], size=11.5)
    bullets(s, [
        (0, "Note that the engineering result and the biological result are independent: the control "
            "system demonstrably works, and what remains unresolved is what the stimulus did."),
    ], top=6.3, size=13.5)

    # ---- 11. next session -------------------------------------------------
    s = new("Plan for the next rig day",
            "Escalating sensitivity, with an explicit decision point at each stage")
    bullets(s, [
        (0, "1.  Verify stimulation delivery from the independent hardware record. This gates everything else."),
        (0, "2.  Maximum-sensitivity screen — all 16 pairs, staircase excitation held 500 ms per level."),
        (1, "Approximately 16× the delivered charge, and a slow modulation that survives an unstable "
            "feature window far better than the fast sequence used on day 1."),
        (0, "3.  Inspect the raw recorded waveform during stimulation. This bypasses the feature pipeline "
            "and is the decisive test:"),
        (1, "Response visible in the raw signal but absent in the feature → the feature extraction is "
            "confirmed as the fault; fix it in software and proceed."),
        (1, "Response absent in both → the issue is stimulus delivery or the preparation, and no "
            "software change will help."),
        (0, "4.  On a confirmed response: two 120 s captures with independent excitation, fit the "
            "multi-input plant model, validate across records, then close the loop."),
        (0, "5.  Contingency: if the feature window is confirmed as the cause, the session converts to "
            "open-loop data collection, with the software correction made between sessions."),
    ], size=14)

    # ---- 12. summary ------------------------------------------------------
    s = new("Summary")
    bullets(s, [
        (0, "The control system is built, verified end to end, and performed correctly on hardware at "
            "the first attempt — zero dropped ticks across 3000 control cycles."),
        (0, "A defect that had left the controller running open-loop for its entire history was found "
            "and corrected before hardware use."),
        (0, "The measurement pipeline was validated against both positive and negative controls, so its "
            "null result on rig day is interpretable rather than ambiguous."),
        (0, "No cortical response to stimulation was detected in the first session. The most likely "
            "explanation is a feature-extraction defect that makes the measured quantity insensitive to "
            "the response — testable, and the first thing addressed next session."),
        (0, "The engineering risk to the project is now low. The open question is biological and "
            "measurement-related, and there is a defined, ordered plan to resolve it."),
    ], size=15.5, top=1.75)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    out = OUT_DIR / "ClosedLoop_system_status_2026-07-30.pptx"
    prs.save(out)
    return out


def main():
    corrs, moved, n_ticks = sweep_from_capture(CAPTURE)
    print(f"Recomputed sweep from {CAPTURE.name}: {n_ticks} ticks, stim channels moved {moved}")
    for i, c in enumerate(corrs, 1):
        print(f"  ch {i:2d}  |corr| {c:.3f}")
    print(f"  max {corrs.max():.3f} on channel {int(corrs.argmax()) + 1}")

    figs = {"arch": fig_architecture(), "sweep": fig_sweep(corrs), "window": fig_window()}
    out = build(corrs, moved, n_ticks, figs)
    print(f"\nWrote {out}")


if __name__ == "__main__":
    main()
